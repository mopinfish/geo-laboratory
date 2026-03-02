"""
exp001 静岡県3都市の道路ネットワーク比較 — 発表資料生成スクリプト

レポート docs/reports/exp001_shizuoka_road_network_report.md の内容を
一般向け10分プレゼンテーション（12枚）として pptx に変換する。

出力: docs/presentations/exp001_shizuoka_road_network_presentation.pptx
"""

from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# --- パス設定 ---
ROOT = Path(__file__).resolve().parent.parent
RESULTS = ROOT / "docs" / "results"
OUTPUT = ROOT / "docs" / "presentations" / "exp001_shizuoka_road_network_presentation.pptx"

# --- 配色 ---
BLUE = RGBColor(0x2B, 0x5C, 0x8A)
RED = RGBColor(0xC0, 0x39, 0x2B)
GRAY = RGBColor(0x55, 0x55, 0x55)
LIGHT_GRAY = RGBColor(0xCC, 0xCC, 0xCC)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x00, 0x00, 0x00)
BG_LIGHT = RGBColor(0xF5, 0xF7, 0xFA)

# フォント
FONT_NAME = "メイリオ"

# スライドサイズ 16:9
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


# =============================================================================
# ユーティリティ
# =============================================================================


def set_font(run, size, color=GRAY, bold=False, name=FONT_NAME):
    """run のフォントを設定"""
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.name = name


def add_textbox(slide, left, top, width, height):
    """テキストボックスを追加して返す"""
    return slide.shapes.add_textbox(left, top, width, height)


def set_paragraph(tf, text, size, color=GRAY, bold=False, alignment=PP_ALIGN.LEFT, spacing_after=Pt(6)):
    """text_frame の最初の段落にテキストを設定"""
    p = tf.paragraphs[0]
    p.alignment = alignment
    p.space_after = spacing_after
    run = p.add_run()
    run.text = text
    set_font(run, size, color, bold)
    return p


def add_paragraph(tf, text, size, color=GRAY, bold=False, alignment=PP_ALIGN.LEFT, spacing_after=Pt(6)):
    """text_frame に段落を追加"""
    p = tf.add_paragraph()
    p.alignment = alignment
    p.space_after = spacing_after
    run = p.add_run()
    run.text = text
    set_font(run, size, color, bold)
    return p


def add_bullet(tf, text, size, color=GRAY, bold=False, level=0, spacing_after=Pt(4)):
    """箇条書きを追加"""
    p = tf.add_paragraph()
    p.level = level
    p.space_after = spacing_after
    run = p.add_run()
    run.text = text
    set_font(run, size, color, bold)
    return p


def fill_background(slide, color):
    """スライド背景を単色で塗りつぶす"""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_bottom_bar(slide, text=""):
    """スライド下部にアクセントバーを追加"""
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), SLIDE_H - Inches(0.35),
        SLIDE_W, Inches(0.35),
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = BLUE
    bar.line.fill.background()
    if text:
        tf = bar.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.RIGHT
        run = p.add_run()
        run.text = text
        set_font(run, 10, WHITE)


def add_slide_number(slide, num, total=12):
    """右下にスライド番号を表示"""
    tb = add_textbox(slide, SLIDE_W - Inches(1.2), SLIDE_H - Inches(0.6), Inches(1.0), Inches(0.4))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    run = p.add_run()
    run.text = f"{num} / {total}"
    set_font(run, 11, LIGHT_GRAY)


def add_section_title(slide, title, subtitle=""):
    """スライド上部にタイトルを配置"""
    # タイトル背景バー
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        SLIDE_W, Inches(1.15),
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = BLUE
    bar.line.fill.background()

    tb = add_textbox(slide, Inches(0.8), Inches(0.15), Inches(11.5), Inches(0.7))
    tf = tb.text_frame
    set_paragraph(tf, title, 30, WHITE, bold=True, alignment=PP_ALIGN.LEFT)
    if subtitle:
        add_paragraph(tf, subtitle, 16, RGBColor(0xBB, 0xD4, 0xEE), alignment=PP_ALIGN.LEFT)


def add_image_safe(slide, img_path, left, top, width=None, height=None):
    """画像を追加（存在確認付き）"""
    p = Path(img_path)
    if not p.exists():
        # プレースホルダーテキスト
        tb = add_textbox(slide, left, top, width or Inches(4), height or Inches(3))
        tf = tb.text_frame
        set_paragraph(tf, f"[画像なし: {p.name}]", 14, RED)
        return None
    if width and height:
        return slide.shapes.add_picture(str(p), left, top, width, height)
    elif width:
        return slide.shapes.add_picture(str(p), left, top, width=width)
    elif height:
        return slide.shapes.add_picture(str(p), left, top, height=height)
    else:
        return slide.shapes.add_picture(str(p), left, top)


def add_rounded_box(slide, left, top, width, height, fill_color=BG_LIGHT, line_color=None):
    """角丸四角形を追加"""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        left, top, width, height,
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape


# =============================================================================
# 各スライドの作成
# =============================================================================


def slide_01_title(prs):
    """スライド1: 表紙"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    fill_background(slide, BLUE)

    # メインタイトル
    tb = add_textbox(slide, Inches(1.5), Inches(1.8), Inches(10.3), Inches(2.0))
    tf = tb.text_frame
    tf.word_wrap = True
    set_paragraph(tf, '道路の"かたち"で街がわかる', 40, WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_paragraph(tf, "静岡県3都市の道路ネットワーク比較", 28, RGBColor(0xBB, 0xD4, 0xEE), alignment=PP_ALIGN.CENTER)

    # 区切り線
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(4.5), Inches(4.2), Inches(4.3), Pt(2),
    )
    line.fill.solid()
    line.fill.fore_color.rgb = WHITE
    line.line.fill.background()

    # サブ情報
    tb2 = add_textbox(slide, Inches(1.5), Inches(4.6), Inches(10.3), Inches(1.5))
    tf2 = tb2.text_frame
    tf2.word_wrap = True
    set_paragraph(tf2, "焼津市 ・ 静岡市 ・ 浜松市", 22, RGBColor(0xDD, 0xDD, 0xDD), alignment=PP_ALIGN.CENTER)
    add_paragraph(tf2, "グラフ理論による道路ネットワーク構造の定量比較", 16, RGBColor(0xAA, 0xBB, 0xCC), alignment=PP_ALIGN.CENTER)
    add_paragraph(tf2, "2026年3月", 14, RGBColor(0x99, 0xAA, 0xBB), alignment=PP_ALIGN.CENTER)


def slide_02_question(prs):
    """スライド2: 問いかけ"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    fill_background(slide, WHITE)
    add_section_title(slide, "道の形は街ごとに違う？")
    add_slide_number(slide, 2)

    # 問いかけ文
    tb = add_textbox(slide, Inches(0.8), Inches(1.5), Inches(11.5), Inches(0.8))
    tf = tb.text_frame
    tf.word_wrap = True
    set_paragraph(tf, "同じ静岡県内でも、歴史や地形が異なれば道路の形も違うはず。", 22, GRAY)
    add_paragraph(tf, "3つの街の道路ネットワークを数値で比べてみました。", 22, GRAY)

    # 3都市カード
    cities = [
        ("焼津市", "漁港の街", "駿河湾に面した漁港都市。\n旧来の漁村集落に由来する\n細い路地が残る。", RGBColor(0xE8, 0xF0, 0xF8)),
        ("静岡市", "城下町", "駿府城の城下町を起源とする\n県庁所在地。格子状の\n市街地と広大な山間部。", RGBColor(0xFD, 0xF0, 0xEC)),
        ("浜松市", "工業都市", "政令指定都市。近代的な\n格子状街路を基盤とする\n遠州灘沿いの工業都市。", RGBColor(0xEE, 0xF5, 0xEE)),
    ]

    card_w = Inches(3.5)
    card_h = Inches(3.5)
    start_x = Inches(0.9)
    gap = Inches(0.55)
    y = Inches(3.0)

    for i, (name, subtitle, desc, bg_color) in enumerate(cities):
        x = start_x + i * (card_w + gap)
        box = add_rounded_box(slide, x, y, card_w, card_h, bg_color, LIGHT_GRAY)

        # 都市名
        tb_name = add_textbox(slide, x + Inches(0.3), y + Inches(0.25), card_w - Inches(0.6), Inches(0.55))
        tf_name = tb_name.text_frame
        set_paragraph(tf_name, name, 26, BLUE, bold=True, alignment=PP_ALIGN.CENTER)

        # サブタイトル
        tb_sub = add_textbox(slide, x + Inches(0.3), y + Inches(0.8), card_w - Inches(0.6), Inches(0.4))
        tf_sub = tb_sub.text_frame
        set_paragraph(tf_sub, subtitle, 18, RED, bold=True, alignment=PP_ALIGN.CENTER)

        # 説明
        tb_desc = add_textbox(slide, x + Inches(0.3), y + Inches(1.4), card_w - Inches(0.6), Inches(1.8))
        tf_desc = tb_desc.text_frame
        tf_desc.word_wrap = True
        set_paragraph(tf_desc, desc, 16, GRAY, alignment=PP_ALIGN.CENTER)


def slide_03_perspective(prs):
    """スライド3: 分析の視点（3つの物差し）"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    fill_background(slide, WHITE)
    add_section_title(slide, "3つの物差しで道路を測る")
    add_slide_number(slide, 3)

    metrics = [
        ("回遊性", "ぐるぐる歩ける？", "道路がループを\nどれだけ形成しているか。\n多いほど迷い歩きを\n楽しめる。", BLUE),
        ("アクセス性", "どこでも行ける？", "交差点の密度や\nネットワーク上の到達しやすさ。\n高いほど移動の\n選択肢が多い。", RGBColor(0x27, 0xAE, 0x60)),
        ("迂回性", "遠回りする？", "最短経路が直線距離の\n何倍になるか。\n低いほど効率的に\n移動できる。", RED),
    ]

    card_w = Inches(3.5)
    card_h = Inches(4.2)
    start_x = Inches(0.9)
    gap = Inches(0.55)
    y = Inches(2.0)

    for i, (title, question, desc, accent) in enumerate(metrics):
        x = start_x + i * (card_w + gap)

        # カード背景
        add_rounded_box(slide, x, y, card_w, card_h, BG_LIGHT, LIGHT_GRAY)

        # アクセントバー上部
        bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            x + Inches(0.6), y + Inches(0.3), card_w - Inches(1.2), Pt(4),
        )
        bar.fill.solid()
        bar.fill.fore_color.rgb = accent
        bar.line.fill.background()

        # タイトル
        tb_t = add_textbox(slide, x + Inches(0.3), y + Inches(0.5), card_w - Inches(0.6), Inches(0.55))
        tf_t = tb_t.text_frame
        set_paragraph(tf_t, title, 26, accent, bold=True, alignment=PP_ALIGN.CENTER)

        # 問い
        tb_q = add_textbox(slide, x + Inches(0.3), y + Inches(1.1), card_w - Inches(0.6), Inches(0.45))
        tf_q = tb_q.text_frame
        set_paragraph(tf_q, question, 20, GRAY, bold=True, alignment=PP_ALIGN.CENTER)

        # 説明
        tb_d = add_textbox(slide, x + Inches(0.3), y + Inches(1.7), card_w - Inches(0.6), Inches(2.0))
        tf_d = tb_d.text_frame
        tf_d.word_wrap = True
        set_paragraph(tf_d, desc, 16, GRAY, alignment=PP_ALIGN.CENTER)

    # 出典注釈
    tb_note = add_textbox(slide, Inches(0.8), Inches(6.6), Inches(11.5), Inches(0.5))
    tf_note = tb_note.text_frame
    set_paragraph(tf_note, "※ 劉（2016）の3観点に基づく道路ネットワーク評価フレームワーク", 12, LIGHT_GRAY, alignment=PP_ALIGN.LEFT)


def slide_04_city_comparison(prs):
    """スライド4: 市全域比較 — レーダーチャート"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    fill_background(slide, WHITE)
    add_section_title(slide, '小さな焼津市が一番"つながっている"', "市全域の道路ネットワーク比較")
    add_slide_number(slide, 4)

    # キーメッセージ
    tb = add_textbox(slide, Inches(0.8), Inches(1.4), Inches(5.5), Inches(1.2))
    tf = tb.text_frame
    tf.word_wrap = True
    set_paragraph(tf, "コンパクトな焼津市が回遊性・アクセス性で\n3市中トップ。", 20, GRAY, bold=True)
    add_paragraph(tf, "広大な山間部を含む静岡市・浜松市は\n市全域でみると数値が低くなる。", 18, GRAY)

    # ポイント
    tb2 = add_textbox(slide, Inches(0.8), Inches(3.0), Inches(5.5), Inches(3.5))
    tf2 = tb2.text_frame
    tf2.word_wrap = True
    set_paragraph(tf2, "注目ポイント", 18, BLUE, bold=True)
    add_bullet(tf2, "焼津市の交差点密度は静岡市の約5倍", 17, GRAY)
    add_bullet(tf2, "静岡市は迂回性が最も高い（山間部の影響）", 17, GRAY)
    add_bullet(tf2, "市域面積: 焼津70km² vs 静岡1,412km² vs 浜松1,558km²", 15, LIGHT_GRAY)

    # レーダーチャート画像
    add_image_safe(slide, RESULTS / "exp001_city_radar_chart.png",
                   Inches(6.5), Inches(1.4), height=Inches(5.5))


def slide_05_coastal(prs):
    """スライド5: 沿岸 vs 内陸"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    fill_background(slide, WHITE)
    add_section_title(slide, "海沿いの方が道が充実していた", "沿岸部 vs 内陸部の比較（仮説を覆す結果）")
    add_slide_number(slide, 5)

    # 仮説と結果
    # 仮説ボックス
    box_h = add_rounded_box(slide, Inches(0.8), Inches(1.5), Inches(5.5), Inches(1.3),
                            RGBColor(0xFD, 0xF0, 0xEC), RED)
    tb_h = add_textbox(slide, Inches(1.0), Inches(1.55), Inches(5.1), Inches(1.2))
    tf_h = tb_h.text_frame
    tf_h.word_wrap = True
    set_paragraph(tf_h, "当初の仮説:", 16, RED, bold=True)
    add_paragraph(tf_h, "「沿岸部は回遊性が低く、迂回性が高い」", 18, RED, bold=True)

    # 結果ボックス
    box_r = add_rounded_box(slide, Inches(0.8), Inches(3.1), Inches(5.5), Inches(1.3),
                            RGBColor(0xE8, 0xF5, 0xE8), RGBColor(0x27, 0xAE, 0x60))
    tb_r = add_textbox(slide, Inches(1.0), Inches(3.15), Inches(5.1), Inches(1.2))
    tf_r = tb_r.text_frame
    tf_r.word_wrap = True
    set_paragraph(tf_r, "実際の結果:", 16, RGBColor(0x27, 0xAE, 0x60), bold=True)
    add_paragraph(tf_r, "3市すべてで沿岸部の回遊性が内陸部より高い！", 18, RGBColor(0x27, 0xAE, 0x60), bold=True)

    # 理由
    tb2 = add_textbox(slide, Inches(0.8), Inches(4.7), Inches(5.5), Inches(2.2))
    tf2 = tb2.text_frame
    tf2.word_wrap = True
    set_paragraph(tf2, "なぜ？", 18, BLUE, bold=True)
    add_bullet(tf2, "沿岸平野部に市街地が集中 → 格子状の道路網が発達", 16, GRAY)
    add_bullet(tf2, "内陸の山間部は谷筋に沿った一本道（樹枝状）が多い", 16, GRAY)
    add_bullet(tf2, "浜松市では沿岸部の交差点密度が内陸部の4.6倍", 16, GRAY)

    # 焼津市沿岸部画像
    add_image_safe(slide, RESULTS / "exp001_焼津市_coastal_screenshot.png",
                   Inches(6.8), Inches(1.4), height=Inches(5.3))


def slide_06_coastal_maps(prs):
    """スライド6: 沿岸部地図比較"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    fill_background(slide, WHITE)
    add_section_title(slide, "3市の沿岸部ネットワーク", "海岸線から2km以内の道路ネットワーク地図")
    add_slide_number(slide, 6)

    img_h = Inches(4.5)
    img_w = Inches(3.7)
    gap = Inches(0.35)
    start_x = Inches(0.6)
    y = Inches(1.5)

    images = [
        (RESULTS / "exp001_焼津市_coastal_screenshot.png", "焼津市", "回遊性 最高 (0.286)"),
        (RESULTS / "exp001_静岡市_coastal_screenshot.png", "静岡市", "迂回性 最高 (1.306)"),
        (RESULTS / "exp001_浜松市_coastal_screenshot.png", "浜松市", "迂回性 最低 (1.118)"),
    ]

    for i, (img_path, label, note) in enumerate(images):
        x = start_x + i * (img_w + gap)
        add_image_safe(slide, img_path, x, y, width=img_w, height=img_h)

        # ラベル
        tb_l = add_textbox(slide, x, y + img_h + Inches(0.1), img_w, Inches(0.4))
        tf_l = tb_l.text_frame
        set_paragraph(tf_l, label, 20, BLUE, bold=True, alignment=PP_ALIGN.CENTER)

        # 注釈
        tb_n = add_textbox(slide, x, y + img_h + Inches(0.45), img_w, Inches(0.35))
        tf_n = tb_n.text_frame
        set_paragraph(tf_n, note, 14, GRAY, alignment=PP_ALIGN.CENTER)


def slide_07_shizuoka_warning(prs):
    """スライド7: 静岡市沿岸部の警告"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    fill_background(slide, WHITE)
    add_section_title(slide, '静岡市の海沿いは"遠回り"が最も多い', "防災上の懸念")
    add_slide_number(slide, 7)

    # 警告ボックス
    add_rounded_box(slide, Inches(0.8), Inches(1.5), Inches(5.5), Inches(1.5),
                    RGBColor(0xFD, 0xF0, 0xEC), RED)
    tb_w = add_textbox(slide, Inches(1.0), Inches(1.55), Inches(5.1), Inches(1.4))
    tf_w = tb_w.text_frame
    tf_w.word_wrap = True
    set_paragraph(tf_w, "静岡市沿岸部の迂回率: 1.306", 22, RED, bold=True)
    add_paragraph(tf_w, "最短経路でも直線距離の約1.3倍を迂回する。", 18, RED)
    add_paragraph(tf_w, "津波避難時に距離が延びるリスク。", 18, RED)

    # 原因と比較
    tb2 = add_textbox(slide, Inches(0.8), Inches(3.3), Inches(5.5), Inches(3.5))
    tf2 = tb2.text_frame
    tf2.word_wrap = True
    set_paragraph(tf2, "原因: 日本平・有度山の丘陵地", 18, BLUE, bold=True)
    add_paragraph(tf2, "清水区と駿河区の間に位置する丘陵地が道路ネットワークを\n南北に分断している。", 16, GRAY)
    add_paragraph(tf2, "", 10, GRAY)
    add_paragraph(tf2, "3市の沿岸部 迂回率の比較:", 16, BLUE, bold=True)
    add_bullet(tf2, "浜松市: 1.118（最も低い = 直線的）", 16, RGBColor(0x27, 0xAE, 0x60))
    add_bullet(tf2, "焼津市: 1.199", 16, GRAY)
    add_bullet(tf2, "静岡市: 1.306（最も高い = 遠回り）", 16, RED, bold=True)

    # 画像
    add_image_safe(slide, RESULTS / "exp001_静岡市_coastal_screenshot.png",
                   Inches(6.8), Inches(1.4), height=Inches(5.3))


def slide_08_station_walk(prs):
    """スライド8: 駅前歩きやすさ"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    fill_background(slide, WHITE)
    add_section_title(slide, '焼津駅が一番"歩き回りやすい"', "駅周辺800mの歩行者ネットワーク比較")
    add_slide_number(slide, 8)

    img_h = Inches(3.5)
    img_w = Inches(3.7)
    gap = Inches(0.35)
    start_x = Inches(0.6)
    y = Inches(1.45)

    images = [
        (RESULTS / "exp001_焼津駅_walk_screenshot.png", "焼津駅", "alpha = 0.330"),
        (RESULTS / "exp001_静岡駅_walk_screenshot.png", "静岡駅", "alpha = 0.308"),
        (RESULTS / "exp001_浜松駅_walk_screenshot.png", "浜松駅", "alpha = 0.247"),
    ]

    for i, (img_path, label, note) in enumerate(images):
        x = start_x + i * (img_w + gap)
        add_image_safe(slide, img_path, x, y, width=img_w, height=img_h)
        tb_l = add_textbox(slide, x, y + img_h + Inches(0.05), img_w, Inches(0.35))
        tf_l = tb_l.text_frame
        set_paragraph(tf_l, label, 18, BLUE, bold=True, alignment=PP_ALIGN.CENTER)
        tb_n = add_textbox(slide, x, y + img_h + Inches(0.35), img_w, Inches(0.3))
        tf_n = tb_n.text_frame
        set_paragraph(tf_n, note, 14, GRAY, alignment=PP_ALIGN.CENTER)

    # 下部コメント
    tb_c = add_textbox(slide, Inches(0.8), Inches(5.7), Inches(11.5), Inches(1.3))
    tf_c = tb_c.text_frame
    tf_c.word_wrap = True
    set_paragraph(tf_c, "意外にも地方都市の焼津駅が最高の回遊性。", 20, BLUE, bold=True)
    add_paragraph(tf_c, "旧漁村集落の細かい路地が、歩行者にとって多様なルート選択を可能にしている。", 17, GRAY)
    add_paragraph(tf_c, "浜松駅は交差点密度こそ最高だが、大きな街区が回遊性を下げている。", 17, GRAY)


def slide_09_alley_power(prs):
    """スライド9: 路地の力"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    fill_background(slide, WHITE)
    add_section_title(slide, "古い路地が街の回遊性を高める", "細街路効果: walk と drive の差")
    add_slide_number(slide, 9)

    # 説明
    tb = add_textbox(slide, Inches(0.8), Inches(1.4), Inches(5.8), Inches(1.5))
    tf = tb.text_frame
    tf.word_wrap = True
    set_paragraph(tf, "歩行者用の細い道（路地・小路）を加えると\n回遊性はどう変わる？", 20, GRAY, bold=True)
    add_paragraph(tf, 'walk（歩行者）と drive（自動車）のネットワークの差で\n"細街路効果"を測定。', 17, GRAY)

    # 3駅の比較ボックス
    box_data = [
        ("焼津駅", "+0.013", "路地が循環路を形成\n→ 回遊性UP", RGBColor(0xE8, 0xF5, 0xE8), RGBColor(0x27, 0xAE, 0x60)),
        ("静岡駅", "-0.005", "ほぼ中立\n城下町の横丁が\n一部寄与", BG_LIGHT, GRAY),
        ("浜松駅", "-0.029", "行き止まりの路地が多い\n→ 回遊性DOWN", RGBColor(0xFD, 0xF0, 0xEC), RED),
    ]

    box_w = Inches(3.5)
    box_h = Inches(2.4)
    start_x = Inches(0.9)
    gap = Inches(0.55)
    y_box = Inches(3.3)

    for i, (name, delta, desc, bg, accent) in enumerate(box_data):
        x = start_x + i * (box_w + gap)
        add_rounded_box(slide, x, y_box, box_w, box_h, bg, accent)

        tb_n = add_textbox(slide, x + Inches(0.2), y_box + Inches(0.15), box_w - Inches(0.4), Inches(0.4))
        tf_n = tb_n.text_frame
        set_paragraph(tf_n, name, 20, accent, bold=True, alignment=PP_ALIGN.CENTER)

        tb_d = add_textbox(slide, x + Inches(0.2), y_box + Inches(0.55), box_w - Inches(0.4), Inches(0.5))
        tf_d = tb_d.text_frame
        set_paragraph(tf_d, f"alpha差分: {delta}", 22, accent, bold=True, alignment=PP_ALIGN.CENTER)

        tb_desc = add_textbox(slide, x + Inches(0.2), y_box + Inches(1.2), box_w - Inches(0.4), Inches(1.0))
        tf_desc = tb_desc.text_frame
        tf_desc.word_wrap = True
        set_paragraph(tf_desc, desc, 14, GRAY, alignment=PP_ALIGN.CENTER)

    # 注釈
    tb_note = add_textbox(slide, Inches(0.8), Inches(6.2), Inches(11.5), Inches(0.8))
    tf_note = tb_note.text_frame
    tf_note.word_wrap = True
    set_paragraph(tf_note, "alpha差分 = walk時のalpha_index − drive時のalpha_index。プラスは細街路が循環路を形成していることを示す。", 12, LIGHT_GRAY)


def slide_10_profiles(prs):
    """スライド10: 3都市の個性"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    fill_background(slide, WHITE)
    add_section_title(slide, "3都市の個性", "ネットワーク特性プロファイル")
    add_slide_number(slide, 10)

    # 左: テキスト要約
    profiles = [
        ("焼津市", "コンパクト高回遊型", "小さな市域に道が密集。路地が\n回遊性を高め、迂回も少ない。", BLUE),
        ("静岡市", "広域高迂回型", "城下町の中心部は歩きやすいが、\n広大な山間部が市全域の数値を下げる。\n沿岸部の遠回りに防災上の課題。", RED),
        ("浜松市", "近代格子・大街区型", "交差点は密だが大きな街区で\n回遊性は低い。\n沿岸部は直線的で避難に有利。", RGBColor(0x27, 0xAE, 0x60)),
    ]

    y_start = Inches(1.5)
    for i, (name, profile_type, desc, accent) in enumerate(profiles):
        y = y_start + i * Inches(1.75)

        # 都市名 + タイプ
        tb_n = add_textbox(slide, Inches(0.8), y, Inches(5.5), Inches(0.5))
        tf_n = tb_n.text_frame
        p = tf_n.paragraphs[0]
        run_name = p.add_run()
        run_name.text = f"{name}  "
        set_font(run_name, 22, accent, bold=True)
        run_type = p.add_run()
        run_type.text = f"— {profile_type}"
        set_font(run_type, 18, GRAY)

        # 説明
        tb_d = add_textbox(slide, Inches(1.2), y + Inches(0.5), Inches(5.1), Inches(1.1))
        tf_d = tb_d.text_frame
        tf_d.word_wrap = True
        set_paragraph(tf_d, desc, 16, GRAY)

    # 右: 棒グラフ画像
    add_image_safe(slide, RESULTS / "exp001_city_comparison_bars.png",
                   Inches(6.8), Inches(1.4), height=Inches(5.5))


def slide_11_summary(prs):
    """スライド11: まとめ"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    fill_background(slide, WHITE)
    add_section_title(slide, "まとめ — 4つの発見")
    add_slide_number(slide, 11)

    findings = [
        ("1", '沿岸部は意外と"つながっている"',
         "3市すべてで沿岸部の回遊性が内陸部より高い。沿岸平野部に市街地が集中しているため。"),
        ("2", '静岡市の海沿いは"遠回り"が多い',
         "日本平・有度山の丘陵地が道路を分断。迂回率1.306は3市の沿岸部で最も高く、津波避難時の懸念。"),
        ("3", '焼津駅は"歩き回りやすさ"No.1',
         "旧漁村の路地が循環路を形成し、回遊性指標は3駅中トップ。地方都市の隠れた資産。"),
        ("4", "古い路地こそ街の宝",
         '焼津は細街路が回遊性を高め、浜松は行き止まりが増える。道路の"質"は都市の歴史で決まる。'),
    ]

    y_start = Inches(1.6)
    for i, (num, title, desc) in enumerate(findings):
        y = y_start + i * Inches(1.4)

        # 番号サークル
        circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(0.8), y, Inches(0.55), Inches(0.55),
        )
        circle.fill.solid()
        circle.fill.fore_color.rgb = BLUE
        circle.line.fill.background()
        tf_c = circle.text_frame
        tf_c.paragraphs[0].alignment = PP_ALIGN.CENTER
        run_c = tf_c.paragraphs[0].add_run()
        run_c.text = num
        set_font(run_c, 20, WHITE, bold=True)

        # タイトル
        tb_t = add_textbox(slide, Inches(1.6), y, Inches(10.5), Inches(0.45))
        tf_t = tb_t.text_frame
        set_paragraph(tf_t, title, 22, BLUE, bold=True)

        # 説明
        tb_d = add_textbox(slide, Inches(1.6), y + Inches(0.45), Inches(10.5), Inches(0.8))
        tf_d = tb_d.text_frame
        tf_d.word_wrap = True
        set_paragraph(tf_d, desc, 16, GRAY)


def slide_12_end(prs):
    """スライド12: 終了"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    fill_background(slide, BLUE)

    # メインテキスト
    tb = add_textbox(slide, Inches(1.5), Inches(2.5), Inches(10.3), Inches(1.5))
    tf = tb.text_frame
    set_paragraph(tf, "ご清聴ありがとうございました", 36, WHITE, bold=True, alignment=PP_ALIGN.CENTER)

    # 区切り線
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(5.0), Inches(4.2), Inches(3.3), Pt(2),
    )
    line.fill.solid()
    line.fill.fore_color.rgb = WHITE
    line.line.fill.background()

    # サブ情報
    tb2 = add_textbox(slide, Inches(1.5), Inches(4.6), Inches(10.3), Inches(2.0))
    tf2 = tb2.text_frame
    tf2.word_wrap = True
    set_paragraph(tf2, '道路の"かたち"で街がわかる', 22, RGBColor(0xBB, 0xD4, 0xEE), alignment=PP_ALIGN.CENTER)
    add_paragraph(tf2, "静岡県3都市の道路ネットワーク比較", 18, RGBColor(0xAA, 0xBB, 0xCC), alignment=PP_ALIGN.CENTER)
    add_paragraph(tf2, "", 10, WHITE, alignment=PP_ALIGN.CENTER)
    add_paragraph(tf2, "データ: OpenStreetMap / 分析: OSMnx + NetworkX", 14, RGBColor(0x88, 0x99, 0xAA), alignment=PP_ALIGN.CENTER)


# =============================================================================
# メイン
# =============================================================================


def main():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    slide_01_title(prs)
    slide_02_question(prs)
    slide_03_perspective(prs)
    slide_04_city_comparison(prs)
    slide_05_coastal(prs)
    slide_06_coastal_maps(prs)
    slide_07_shizuoka_warning(prs)
    slide_08_station_walk(prs)
    slide_09_alley_power(prs)
    slide_10_profiles(prs)
    slide_11_summary(prs)
    slide_12_end(prs)

    OUTPUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUTPUT))
    print(f"生成完了: {OUTPUT}")
    print(f"スライド数: {len(prs.slides)}")


if __name__ == "__main__":
    main()
