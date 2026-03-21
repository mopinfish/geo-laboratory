"""exp002 プレゼンテーション生成スクリプト"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
import os

# --- 定数 ---
SLIDE_W = Emu(12192000)
SLIDE_H = Emu(6858000)
TOTAL_SLIDES = 12
PHOTO_DIR = os.path.join("docs", "results", "exp002", "photos")

# 色定義
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x1A, 0x1A, 0x2E)
DARK_BLUE = RGBColor(0x00, 0x33, 0x66)
BLUE = RGBColor(0x00, 0x66, 0xCC)
LIGHT_BLUE = RGBColor(0xE8, 0xF4, 0xFD)
ACCENT = RGBColor(0x00, 0x99, 0xCC)
GRAY = RGBColor(0x66, 0x66, 0x66)
LIGHT_GRAY = RGBColor(0xAA, 0xAA, 0xAA)
GREEN = RGBColor(0x2E, 0x7D, 0x32)
ORANGE = RGBColor(0xE6, 0x7E, 0x22)

OUTPUT_DIR = "docs/presentations"
IMG_DIR = "tmp"


def add_textbox(slide, left, top, width, height, text, font_size=16,
                bold=False, color=BLACK, alignment=PP_ALIGN.LEFT, font_name="Meiryo"):
    """テキストボックスを追加する"""
    txBox = slide.shapes.add_textbox(Emu(left), Emu(top), Emu(width), Emu(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = alignment
    run = p.runs[0]
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font_name
    return txBox


def add_multiline_textbox(slide, left, top, width, height, lines, font_size=16,
                          color=BLACK, alignment=PP_ALIGN.LEFT, line_spacing=1.2,
                          font_name="Meiryo"):
    """複数行テキストボックスを追加する"""
    txBox = slide.shapes.add_textbox(Emu(left), Emu(top), Emu(width), Emu(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, (text, bold) in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = alignment
        p.space_after = Pt(font_size * 0.3)
        if text:
            run = p.add_run()
            run.text = text
            run.font.size = Pt(font_size)
            run.font.bold = bold
            run.font.color.rgb = color
            run.font.name = font_name
    return txBox


def add_page_number(slide, num):
    """ページ番号を追加する"""
    add_textbox(slide, 11200000, 6400000, 800000, 300000,
                f"{num} / {TOTAL_SLIDES}", font_size=11, color=LIGHT_GRAY,
                alignment=PP_ALIGN.RIGHT)


def build_presentation():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    blank_layout = prs.slide_layouts[6]  # Blank

    # ============================================================
    # Slide 1: タイトル
    # ============================================================
    slide = prs.slides.add_slide(blank_layout)
    # 背景色風の帯（上部）
    add_textbox(slide, 1200000, 1200000, 9800000, 1200000,
                "石の島の水を探す", font_size=40, bold=True, color=DARK_BLUE,
                alignment=PP_ALIGN.CENTER)
    add_textbox(slide, 1200000, 2500000, 9800000, 800000,
                "衛星画像で見つける北木島の丁場跡", font_size=28, color=BLUE,
                alignment=PP_ALIGN.CENTER)
    add_textbox(slide, 1200000, 3600000, 9800000, 600000,
                "岡山県笠岡市・北木島  |  Sentinel-2 衛星画像  |  水域指数（NDWI・MNDWI）",
                font_size=16, color=GRAY, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, 1200000, 5800000, 9800000, 400000,
                "2026年3月", font_size=14, color=LIGHT_GRAY,
                alignment=PP_ALIGN.CENTER)

    # ============================================================
    # Slide 2: 北木島ってどんな島？
    # ============================================================
    slide = prs.slides.add_slide(blank_layout)
    add_textbox(slide, 500000, 300000, 11000000, 700000,
                "北木島ってどんな島？", font_size=30, bold=True, color=DARK_BLUE)
    add_page_number(slide, 2)

    add_multiline_textbox(slide, 500000, 1200000, 5500000, 5000000, [
        ("笠岡諸島で一番大きな島", True),
        ("岡山県笠岡市に属し、花崗岩（北木石）の産地として全国に知られています。", False),
        ("", False),
        ("「石の島」としての歴史", True),
        ("1620年 — 大坂城の石垣に北木石が使用", False),
        ("1892年 — 鶴田丁場が開業、産業としての採石が本格化", False),
        ("1957年 — 丁場（採石場）が島内に127か所、人口は最大12,000人", False),
        ("現在  — 丁場はわずか2か所、人口は約600〜700人", False),
        ("", False),
        ("北木石が使われた有名な建物", True),
        ("大坂城石垣 / 日本銀行本店 / 靖国神社大鳥居 / 明治神宮橋 /", False),
        ("三越日本橋本店 / 東京駅 / 五条大橋", False),
    ], font_size=16, color=BLACK)

    add_multiline_textbox(slide, 6500000, 1200000, 5200000, 5000000, [
        ("2019年 日本遺産に認定", True),
        ("「知ってる!? 悠久の時が流れる石の島」", False),
        ("", False),
        ("丁場跡に水が溜まる", True),
        ("閉鎖後の採石場は深い凹地に\n雨水が溜まり、池になっています。\n「北木の桂林」と呼ばれる\nエメラルドグリーンの湖も。", False),
        ("", False),
        ("でも…全部でいくつある？", True),
        ("丁場跡の全容を把握した\n調査は限られています。\nそこで衛星画像の出番です！", False),
    ], font_size=16, color=BLACK)

    # ============================================================
    # Slide 3: 現地の様子（写真スライド）
    # ============================================================
    slide = prs.slides.add_slide(blank_layout)
    add_textbox(slide, 500000, 300000, 11000000, 700000,
                "現地の様子 — 丁場跡と瀬戸内の風景", font_size=30, bold=True, color=DARK_BLUE)
    add_page_number(slide, 3)

    # 4枚の写真を2x2で配置
    photos = [
        (os.path.join(PHOTO_DIR, "choba_lake_1.jpg"), "丁場湖 — 切り立った花崗岩の\n岩壁に囲まれた水域"),
        (os.path.join(PHOTO_DIR, "stage_2.jpg"), "湖上ステージ — 採石跡を\n活用したイベント会場"),
        (os.path.join(PHOTO_DIR, "venice_1.jpg"), "瀬戸内のベニス — 廃石で\n築かれた護岸と透明な水"),
        (os.path.join(PHOTO_DIR, "sea_1.jpg"), "穏やかな瀬戸内の海 —\n北木島の港の風景"),
    ]
    positions = [
        (300000, 1100000),
        (6100000, 1100000),
        (300000, 3800000),
        (6100000, 3800000),
    ]
    img_w = Emu(5600000)
    img_h = Emu(2500000)

    for (photo_path, caption), (px, py) in zip(photos, positions):
        if os.path.exists(photo_path):
            slide.shapes.add_picture(photo_path, Emu(px), Emu(py), img_w, img_h)
        add_textbox(slide, px, py + 2500000, 5600000, 500000,
                    caption, font_size=12, color=GRAY, alignment=PP_ALIGN.CENTER)

    add_textbox(slide, 500000, 6400000, 11000000, 300000,
                "2026年3月 現地訪問時に撮影",
                font_size=11, color=LIGHT_GRAY, alignment=PP_ALIGN.RIGHT)

    # ============================================================
    # Slide 4: やりたいこと
    # ============================================================
    slide = prs.slides.add_slide(blank_layout)
    add_textbox(slide, 500000, 300000, 11000000, 700000,
                "衛星画像で丁場跡の「水」を見つけたい", font_size=30, bold=True, color=DARK_BLUE)
    add_textbox(slide, 500000, 1000000, 11000000, 600000,
                "水は光の特定の波長を吸収する性質があります。この性質を利用して、衛星画像から水域だけを浮かび上がらせます。",
                font_size=20, color=BLACK)
    add_page_number(slide, 4)

    # 3カラムで手法を説明 (slide 4)
    col_w = 3400000
    col_gap = 300000
    col_left = 500000
    col_top = 2000000

    for i, (title, subtitle, body) in enumerate([
        ("NDWI",
         "緑色 vs 近赤外線",
         "水は近赤外線を強く吸収。\n緑色との差を計算すると\n水域がプラスの値になる。\n\nマクフィーターズ（1996）が提唱。"),
        ("MNDWI",
         "緑色 vs 短波長赤外線",
         "近赤外線の代わりに\n短波長赤外線を使う改良版。\n岩肌と水域の区別が\nNDWIよりうまくいく。\n\nシュー（2006）が提唱。"),
        ("合わせ技で検出",
         "2つの指数の「いいとこ取り」",
         "NDWIかMNDWIの\nどちらかで水域と判定された\nピクセルを水域とみなす。\n\nさらに植生（NDVI）が高い\nエリアは除外して精度UP。"),
    ]):
        x = col_left + i * (col_w + col_gap)
        add_textbox(slide, x, col_top, col_w, 500000,
                    title, font_size=26, bold=True, color=BLUE,
                    alignment=PP_ALIGN.CENTER)
        add_textbox(slide, x, col_top + 500000, col_w, 400000,
                    subtitle, font_size=16, bold=True, color=ACCENT,
                    alignment=PP_ALIGN.CENTER)
        add_textbox(slide, x, col_top + 1000000, col_w, 3000000,
                    body, font_size=16, color=BLACK, alignment=PP_ALIGN.CENTER)

    # ============================================================
    # Slide 5: 使ったデータ
    # ============================================================
    slide = prs.slides.add_slide(blank_layout)
    add_textbox(slide, 500000, 300000, 11000000, 700000,
                "使ったデータ", font_size=30, bold=True, color=DARK_BLUE)
    add_page_number(slide, 5)

    add_multiline_textbox(slide, 500000, 1200000, 5500000, 5000000, [
        ("Sentinel-2 衛星", True),
        ("ヨーロッパ宇宙機関（ESA）が運用する\n地球観測衛星。10mの解像度で\n地表を撮影し、データは無料公開。", False),
        ("", False),
        ("Microsoft Planetary Computer", True),
        ("衛星画像をクラウドで検索・取得できる\n無料サービス。プログラムから\n直接アクセスできます。", False),
        ("", False),
        ("使用した光の波長", True),
        ("緑色（560nm）/ 赤色（665nm）/ 近赤外線（842nm）/ 短波長赤外線（1610nm）", False),
    ], font_size=16, color=BLACK)

    add_multiline_textbox(slide, 6500000, 1200000, 5200000, 5000000, [
        ("2つの時期の画像を比較", True),
        ("", False),
        ("春の画像（2025年3月23日）", True),
        ("雲量 0.0%   植生が少ない時期", False),
        ("岩肌と水の区別がやや難しい", False),
        ("", False),
        ("夏の画像（2025年8月2日）", True),
        ("雲量 0.7%   植生が豊富な時期", False),
        ("緑の島 + 水域 が鮮明に分かれる", False),
        ("", False),
        ("→ 夏の画像を中心に分析", True),
    ], font_size=16, color=BLACK)

    # ============================================================
    # Slide 5: 水域指数の地図（4パネル）
    # ============================================================
    slide = prs.slides.add_slide(blank_layout)
    add_textbox(slide, 500000, 300000, 11000000, 700000,
                "衛星画像から計算した3つの指数", font_size=30, bold=True, color=DARK_BLUE)
    add_textbox(slide, 500000, 900000, 11000000, 400000,
                "青い部分が水域、赤い部分が陸地。右下が最終的な判定結果（青=水域、緑=植生、灰=その他）",
                font_size=16, color=GRAY)
    add_page_number(slide, 6)

    img_path = os.path.join(IMG_DIR, "exp002_ndwi_static.png")
    if os.path.exists(img_path):
        slide.shapes.add_picture(img_path,
                                 Emu(800000), Emu(1400000),
                                 Emu(10600000), Emu(5200000))

    # ============================================================
    # Slide 6: ヒストグラム
    # ============================================================
    slide = prs.slides.add_slide(blank_layout)
    add_textbox(slide, 500000, 300000, 11000000, 700000,
                "指数の分布 — 水と陸を分けるライン", font_size=30, bold=True, color=DARK_BLUE)
    add_page_number(slide, 7)

    add_textbox(slide, 500000, 900000, 11000000, 400000,
                "破線が判定の「しきい値」。左側が陸、右側が水域。しきい値を低めに設定して、小さな水域も拾えるようにしています。",
                font_size=16, color=GRAY)

    img_path = os.path.join(IMG_DIR, "exp002_ndwi_histogram.png")
    if os.path.exists(img_path):
        slide.shapes.add_picture(img_path,
                                 Emu(400000), Emu(1600000),
                                 Emu(11400000), Emu(4800000))

    # ============================================================
    # Slide 7: 水域強調画像
    # ============================================================
    slide = prs.slides.add_slide(blank_layout)
    add_textbox(slide, 500000, 300000, 11000000, 700000,
                "衛星写真で見る北木島 — 水域を青く強調", font_size=30, bold=True, color=DARK_BLUE)
    add_page_number(slide, 8)

    add_textbox(slide, 500000, 900000, 5000000, 400000,
                "左: 衛星写真そのまま　　右: 水域を青色で強調",
                font_size=16, color=GRAY)

    img_path = os.path.join(IMG_DIR, "exp002_geotiff_preview.png")
    if os.path.exists(img_path):
        slide.shapes.add_picture(img_path,
                                 Emu(400000), Emu(1400000),
                                 Emu(11400000), Emu(5200000))

    add_textbox(slide, 500000, 6200000, 11000000, 400000,
                "島の内部に点在する青いスポットが、丁場跡に溜まった水域の候補です",
                font_size=16, bold=True, color=BLUE)

    # ============================================================
    # Slide 8: 検出結果の数字
    # ============================================================
    slide = prs.slides.add_slide(blank_layout)
    add_textbox(slide, 500000, 300000, 11000000, 700000,
                "見つかった水域 — 145か所", font_size=30, bold=True, color=DARK_BLUE)
    add_page_number(slide, 9)

    # 大きな数字
    add_textbox(slide, 500000, 1200000, 5000000, 1200000,
                "145", font_size=72, bold=True, color=BLUE,
                alignment=PP_ALIGN.CENTER)
    add_textbox(slide, 500000, 2400000, 5000000, 500000,
                "島内で検出された水域ポリゴン", font_size=20, color=BLACK,
                alignment=PP_ALIGN.CENTER)

    add_multiline_textbox(slide, 500000, 3200000, 5000000, 3000000, [
        ("面積の分布", True),
        ("100〜500m² … 112か所（小さな水溜り）", False),
        ("500〜1,000m² … 19か所", False),
        ("1,000〜5,000m² … 12か所", False),
        ("5,000m²以上 … 2か所（大きな丁場跡）", False),
    ], font_size=16, color=BLACK)

    add_multiline_textbox(slide, 6200000, 1200000, 5500000, 5000000, [
        ("4つの集中エリア", True),
        ("", False),
        ("島の北部（最大の集中地帯）", True),
        ("最大水域 7,826m² を含む。\n上位10件中6件がここに集中。\n金風呂地区の採石中心地に対応？", False),
        ("", False),
        ("島の南東部", True),
        ("2番目に大きい水域 6,521m²。\n「北木の桂林」に相当か。", False),
        ("", False),
        ("島の中央部", True),
        ("中規模の水域が散在。", False),
        ("", False),
        ("島の西部", True),
        ("小〜中規模の水域が点在。", False),
    ], font_size=16, color=BLACK)

    # ============================================================
    # Slide 9: 春 vs 夏
    # ============================================================
    slide = prs.slides.add_slide(blank_layout)
    add_textbox(slide, 500000, 300000, 11000000, 700000,
                "春と夏で見え方が変わる", font_size=30, bold=True, color=DARK_BLUE)
    add_textbox(slide, 500000, 900000, 11000000, 400000,
                "同じ島を撮った2つの季節の衛星画像を比較しました",
                font_size=18, color=GRAY)
    add_page_number(slide, 10)

    # 2カラム比較
    for i, (season, date, items) in enumerate([
        ("春（3月23日）", "植生が少なく岩肌が目立つ", [
            ("検出ポリゴン", "113か所"),
            ("最大水域", "1.28 ha"),
            ("1,000m²超", "22か所"),
            ("NDWI最大値", "1.000"),
            ("特徴", "大きな水域を見つけやすい"),
        ]),
        ("夏（8月2日）", "緑豊かで島の輪郭が鮮明", [
            ("検出ポリゴン", "145か所"),
            ("最大水域", "0.78 ha"),
            ("1,000m²超", "14か所"),
            ("NDWI最大値", "0.210"),
            ("特徴", "小さな水域まで拾える"),
        ]),
    ]):
        x = 500000 + i * 6000000
        add_textbox(slide, x, 1500000, 5500000, 600000,
                    season, font_size=26, bold=True,
                    color=GREEN if i == 1 else ORANGE,
                    alignment=PP_ALIGN.CENTER)
        add_textbox(slide, x, 2100000, 5500000, 400000,
                    date, font_size=16, color=GRAY,
                    alignment=PP_ALIGN.CENTER)

        for j, (label, value) in enumerate(items):
            y = 2700000 + j * 500000
            add_textbox(slide, x, y, 2500000, 400000,
                        label, font_size=16, bold=True, color=BLACK)
            add_textbox(slide, x + 2800000, y, 2700000, 400000,
                        value, font_size=16, color=BLUE if i == 1 else ORANGE)

    add_textbox(slide, 500000, 5700000, 11000000, 800000,
                "→ 季節ごとに得意な検出パターンが違う。両方の結果を合わせると、より網羅的に丁場跡を把握できそうです。",
                font_size=18, bold=True, color=DARK_BLUE)

    # ============================================================
    # Slide 11: まとめ
    # ============================================================
    slide = prs.slides.add_slide(blank_layout)
    add_textbox(slide, 500000, 300000, 11000000, 700000,
                "まとめ — わかったこと", font_size=30, bold=True, color=DARK_BLUE)
    add_page_number(slide, 11)

    findings = [
        ("1", "衛星画像から丁場跡の水域を検出できた",
         "NDWIとMNDWIの合わせ技で、島内に145か所の水域を発見。\n歴史記録の127か所の丁場に近い数です。"),
        ("2", "水域は島の北部に集中している",
         "金風呂地区（採石の中心地）に対応する地域に\n最も多くの水域が見つかりました。"),
        ("3", "春と夏で検出の特性が異なる",
         "春は大型水域、夏は小規模水域の検出に強い。\n両方を使うと網羅性が向上します。"),
        ("4", "次のステップ — 現地で確かめる",
         "衛星からの検出はあくまで「候補」。\n実際の丁場跡との照合が今後の課題です。"),
    ]

    for i, (num, title, body) in enumerate(findings):
        y = 1100000 + i * 1350000
        add_textbox(slide, 500000, y, 600000, 600000,
                    num, font_size=28, bold=True, color=WHITE,
                    alignment=PP_ALIGN.CENTER)
        # 番号の背景円（簡易的にテキストで代用）
        shape = slide.shapes.add_shape(
            1, Emu(500000), Emu(y), Emu(550000), Emu(550000))  # oval
        shape.fill.solid()
        shape.fill.fore_color.rgb = BLUE
        shape.line.fill.background()
        # 番号テキストを再度上に
        add_textbox(slide, 500000, y + 100000, 550000, 400000,
                    num, font_size=24, bold=True, color=WHITE,
                    alignment=PP_ALIGN.CENTER)

        add_textbox(slide, 1300000, y, 10400000, 450000,
                    title, font_size=22, bold=True, color=DARK_BLUE)
        add_textbox(slide, 1300000, y + 450000, 10400000, 800000,
                    body, font_size=16, color=GRAY)

    # ============================================================
    # Slide 12: 終了
    # ============================================================
    slide = prs.slides.add_slide(blank_layout)
    add_textbox(slide, 1200000, 2000000, 9800000, 1200000,
                "ありがとうございました", font_size=36, bold=True, color=DARK_BLUE,
                alignment=PP_ALIGN.CENTER)
    add_textbox(slide, 1200000, 3300000, 9800000, 600000,
                "石の島の水を探す — 衛星画像で見つける北木島の丁場跡",
                font_size=22, color=BLUE, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, 1200000, 4200000, 9800000, 400000,
                "データ: Sentinel-2 (ESA) / Microsoft Planetary Computer",
                font_size=14, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, 1200000, 4600000, 9800000, 400000,
                "分析: Python / rasterio / folium / uv",
                font_size=14, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

    # --- 保存 ---
    os.makedirs(OUTPUT_DIR, exist_ok=True)
    output_path = os.path.join(OUTPUT_DIR,
                               "exp002_kitagi_quarry_water_detection_presentation.pptx")
    prs.save(output_path)
    print(f"プレゼンテーションを保存しました: {output_path}")
    print(f"スライド数: {len(prs.slides)}")


if __name__ == "__main__":
    build_presentation()
