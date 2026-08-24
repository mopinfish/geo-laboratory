#!/usr/bin/env python3
"""FOSS4G 2026 北木島 口頭発表（英語・12枚、または再訪なし版11枚）の PPTX を生成する。

各スライドにタイトル・本文（必要に応じて callout・footer）・図版/写真・スライド
番号を配置する。図版・写真は `docs/presentations/images/` から読み込む
（外部ソースからのコピー元と SHA256 は Task 6 の照合記録に記載）。

タイトル文字列・投影本文は内容契約（docs/presentations/exp002_kitagi_foss4g2026_presentation.md）
の `## Slide N — ` 見出しと `Projected body` から転記した完全一致文字列。

`--no-revisit` を指定すると S9（8/31 再訪スライド）を除いた11枚を
`exp002_kitagi_foss4g2026_presentation_no_revisit.pptx` として出力する
（現地訪問が未実施の場合の切替。`build(revisit=False)` が11枚を返す）。

スピーカーノートは `exp002_kitagi_foss4g2026_presentation_speaker_notes.md` を正本とし、
ビルド時に読み込んで各スライドのノートペインへ書き込む（英語の発話本文を先頭に、
空行を挟んで日本語の非発話部分を続ける）。再訪なし版では S9 の節を使わないだけで、
残りのスライドは内容契約のスライド番号で対応付ける。

決定性（determinism）: 生成される .pptx はコミット対象であり、再生成しても
バイト単位で同一になる必要がある。python-pptx が自動的にスタンプしうる
docProps/core.xml のプロパティ（created / modified / last_modified_by /
revision / author / title）はすべて固定リテラル値で明示的に上書きする。

使い方:
    uv run python docs/presentations/exp002_kitagi_foss4g2026_presentation.py
    uv run python docs/presentations/exp002_kitagi_foss4g2026_presentation.py --no-revisit
"""

from __future__ import annotations

import argparse
import datetime as dt
import io
import re
import zipfile
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Inches, Pt

BASE = Path(__file__).resolve().parent
IMAGES = BASE / "images"
OUTPUT = BASE / "exp002_kitagi_foss4g2026_presentation.pptx"
NO_REVISIT_OUTPUT = BASE / "exp002_kitagi_foss4g2026_presentation_no_revisit.pptx"
NOTES_MD = BASE / "exp002_kitagi_foss4g2026_presentation_speaker_notes.md"

# スピーカーノート Markdown の構造マーカー（英語＝発話対象、日本語＝非発話）。
EN_MARKER = "**EN (spoken)**"
JA_MARKER = "**JA (not spoken)**"

# --- 16:9 スライドサイズ ---
SLIDE_W = Emu(12192000)
SLIDE_H = Emu(6858000)

# --- レイアウト定数 ---
MARGIN = Inches(0.55)
TITLE_TOP = Inches(0.4)
TITLE_HEIGHT = Inches(1.8)
NUM_W = Inches(0.8)
NUM_H = Inches(0.35)

# --- タイプスケール ---
SZ_TITLE = 28
SZ_BODY = 17
SZ_BODY_NARROW = 16  # 図版を実寸近くまで拡大して本文列が狭くなるスライド（S2）用。下限15ptは維持
SZ_CALLOUT = 66  # evidence 階層の主結果（S6 の 145 等）60〜72pt の中央値
SZ_NOTE = 14
SZ_FOOTER = 11  # 帰属・ライセンス・ライブラリ名（11〜12pt）
SZ_NUM = 11
BODY_MIN_PT = 15  # 本文はこの下限を下回ってはならない

# --- 配色 ---
COL_TEXT = RGBColor(0x2B, 0x2B, 0x2B)
COL_ACCENT = RGBColor(0x0D, 0x47, 0xA1)
COL_MUTED = RGBColor(0x5A, 0x56, 0x4E)

# --- 図版・写真配置の共通ジオメトリ ---
# タイトル直下からフッター/スライド番号の手前までの領域（画像・本文を配置する帯）。
GRAPHIC_TOP = TITLE_TOP + TITLE_HEIGHT
GRAPHIC_BOTTOM = SLIDE_H - Inches(0.5)
GRAPHIC_HEIGHT = GRAPHIC_BOTTOM - GRAPHIC_TOP
USABLE_LEFT = MARGIN
USABLE_WIDTH = SLIDE_W - 2 * MARGIN
IMAGE_TEXT_GAP = Inches(0.3)

# 写真スロット（S3・S4・S9）: 固定寸法、16:9、中央クロップして配置する。
# 差し替え（例: 8/31 撮影分の投入）で生成スクリプトを変更しなくてよいよう、
# 位置決めのロジックはこの2定数だけに依存させる。
PHOTO_SLOT_W_IN = Inches(5.6)
PHOTO_SLOT_H_IN = Inches(3.15)
PHOTO_ROW_GAP = Inches(0.2)
PHOTO_ROW_TOP = GRAPHIC_BOTTOM - PHOTO_SLOT_H_IN
PHOTO_ROW_LEFT = USABLE_LEFT + (USABLE_WIDTH - 2 * PHOTO_SLOT_W_IN - PHOTO_ROW_GAP) // 2
# 写真スロットの上に置く本文（S3・S4・S9 共通）の高さ。
PHOTO_TEXT_HEIGHT = PHOTO_ROW_TOP - GRAPHIC_TOP - Inches(0.15)

# S4 の写真スロット: **縦長**写真2枚を同寸で並置する（S3・S9 の16:9固定スロットとは別系統）。
# S4 で使う2枚（`aerial_quarry_pond.jpg` / `drone_lake_stage.jpg`）はいずれも縦位置
# （縦横比 0.75〜0.88）で、16:9 スロットに収めると各フレームの 58% を捨てることになる。
# S4 の主張は「上空から見ると採石権の境界が地形として読める」であり、それには縦方向の
# 広がりが必要なため、スロット自体を縦長にして本文列と左右に並べる。
# スロットの縦横比は `aerial_quarry_pond.jpg`（UI除去後 1080 x 1230 px）の実比に合わせ、
# この写真をクロップなしで全面表示する。もう1枚（1080 x 1440 px）は上下クロップになり、
# 被写体の無い下部（板張りの床）を落とす（`S04_PHOTO_VBIAS_2`）。
S04_PHOTO_SLOT_AR = 1080 / 1230
S04_PHOTO_SLOT_H = Inches(4.4)
S04_PHOTO_GAP = Inches(0.2)
S04_PHOTO_VBIAS_2 = 0.0  # 縦クロップは全量を下側から取る（空の板張りの床を落とす）

# S2 の位置図（`poster_f1_study_area.png`）の配置倍率。
# この図版の最小 native フォントは 15pt（パネル(b)の島名ラベル）なので、
# 実効 15pt 以上には 倍率 ≧ 1.0 が必要。丸め・tight bbox の余白ぶんの余裕を取り 1.04。
# 実寸 7.85 x 4.50 in なので配置は 8.16 x 4.68 in となり、図版帯（高さ4.8in）に収まる。
S02_MAP_SCALE = 1.04


def _image_size(path: Path) -> tuple[int, int]:
    with Image.open(path) as im:
        return im.size


def add_picture_contain(slide, path: Path, left, top, max_width, max_height, name: str):
    """画像をアスペクト比を保ったまま矩形内に収めて中央配置する（クロップしない）。

    地図・グラフ・QRコードなど情報を伝える図版はラベル欠落を避けるためクロップせず、
    矩形内で均一スケールし中央に置く。
    """
    iw, ih = _image_size(path)
    scale = min(max_width / iw, max_height / ih)
    width = int(iw * scale)
    height = int(ih * scale)
    pic_left = int(left + (max_width - width) / 2)
    pic_top = int(top + (max_height - height) / 2)
    pic = slide.shapes.add_picture(str(path), pic_left, pic_top, width=width, height=height)
    pic.name = name
    return pic


def picture_native_size_in(path: Path) -> tuple[float, float]:
    """画像の実寸（インチ）を返す。ピクセル数 ÷ PNG/JPEG に記録された dpi。

    図中に焼き込まれた文字がスライド上で何ptに見えるかは
    `native_pt × (配置幅 ÷ 実寸幅)` で決まるため、実寸は配置設計の基礎量である。
    """
    with Image.open(path) as im:
        px_w, px_h = im.size
        dpi_x, dpi_y = im.info.get("dpi", (72.0, 72.0))
    return px_w / dpi_x, px_h / dpi_y


def add_picture_at_scale(slide, path: Path, left, top, scale: float, name: str):
    """画像を「実寸 × scale」で配置する（矩形に収める contain とは別の指定方法）。

    図中の焼き込み文字の実効ptを直接指定したい図版に使う。実効pt は
    `native_pt × scale` になるので、`scale` を決めれば下限（15pt）の充足が
    レイアウトの偶然に左右されない。

    ポスターから流用した `poster_f1_study_area.png`（S2）はネットワーク取得した
    地理院タイルを含み再生成できないため、native 15pt を実効15pt以上にするには
    ほぼ実寸で置く必要がある（`scale ≥ 1.0`）。その意図をコードに残すため、
    「幅5.4inの箱に収める」ではなく「実寸の1.04倍で置く」と書く。
    """
    w_in, h_in = picture_native_size_in(path)
    width = Emu(int(round(914400 * w_in * scale)))
    height = Emu(int(round(914400 * h_in * scale)))
    pic = slide.shapes.add_picture(str(path), left, top, width=width, height=height)
    pic.name = name
    return pic


def add_picture_cover(slide, path: Path, left, top, width, height, name: str, *, vbias: float = 0.5):
    """画像を矩形いっぱいにクロップで配置する（縦横比は変形させない）。

    写真スロット（S1 表紙・S3・S4・S9）で使用する。矩形とのアスペクト比の差分だけを
    長辺方向からクロップし、残った部分を矩形にそのまま引き伸ばす（この時点で
    残存部分のアスペクト比は矩形と一致しているため、実際には変形しない）。

    縦方向のクロップ（画像が縦長でクロップ対象が上下になる場合）が必要なとき、
    `vbias`（0〜1、既定 0.5 で中央対称）は必要クロップ量のうち上側から取る割合を
    表す。`vbias` を大きくすると上側をより多くクロップして構図が下寄りに動き
    （例: 水面を残す）、小さくすると下側をより多くクロップして構図が上寄りに動く
    （例: 写り込んだ機材・UIを画面下端から外す）。個体差のある画像アスペクト比
    でも必要クロップ量から比率で再計算するため、複数画像へ同じ `vbias` を安全に
    使い回せる。
    """
    iw, ih = _image_size(path)
    img_ar = iw / ih
    target_ar = width / height
    pic = slide.shapes.add_picture(str(path), left, top, width=width, height=height)
    pic.name = name
    if img_ar > target_ar:
        visible = target_ar / img_ar
        crop = (1 - visible) / 2
        pic.crop_left = crop
        pic.crop_right = crop
    elif img_ar < target_ar:
        visible = img_ar / target_ar
        total = 1 - visible
        pic.crop_top = total * vbias
        pic.crop_bottom = total * (1 - vbias)
    return pic


def add_photo_pair(
    slide, path1: Path, path2: Path, name1: str = "Picture1", name2: str = "Picture2",
    *, vbias1: float = 0.5, vbias2: float = 0.5,
):
    """写真スロット2枚を同寸で並置する（S3・S4・S9 共通）。

    `vbias1`/`vbias2` は各写真のクロップ位置（`add_picture_cover` の `vbias`）を
    個別に指定する場合に使う（既定は中央対称クロップ）。
    """
    add_picture_cover(
        slide, path1, PHOTO_ROW_LEFT, PHOTO_ROW_TOP, PHOTO_SLOT_W_IN, PHOTO_SLOT_H_IN, name1,
        vbias=vbias1,
    )
    add_picture_cover(
        slide, path2,
        PHOTO_ROW_LEFT + PHOTO_SLOT_W_IN + PHOTO_ROW_GAP,
        PHOTO_ROW_TOP, PHOTO_SLOT_W_IN, PHOTO_SLOT_H_IN, name2,
        vbias=vbias2,
    )


def resolve_revisit_photo(n: int) -> Path:
    """8/31 再訪の撮影写真を解決する（撮影後の差し替えを生成スクリプト変更なしで行う）。

    `images/revisit_{n}.jpg`（`.jpeg` / `.png` も許容）が存在すればそれを使う。
    撮影がまだの間は `images/placeholder_revisit_{n}.png`
    （`Placeholder — 2026-08-31 photograph` と描いた無地画像）を使う。
    """
    for ext in ("jpg", "jpeg", "png"):
        candidate = IMAGES / f"revisit_{n}.{ext}"
        if candidate.is_file():
            return candidate
    return IMAGES / f"placeholder_revisit_{n}.png"


# 内容契約の `## Slide N — ` 見出しから転記した、各スライドのタイトル完全一致文字列。
TITLES = [
    "Detecting Quarry Pond Remnants on a Japanese Island Heritage Site Using Sentinel-2 Imagery and Open-Source Remote Sensing Tools",
    "A quarrying island, and the ponds it left behind",
    "On foot: I could stand in front of five or six of them",
    "From the air: the quarry boundaries are the landform",
    "On the train home: one satellite scene covers the whole island",
    "The scan found 145 water polygons across the island",
    "Each scale shows what the others cannot",
    "Five or six sites visited — the scan produced 145 candidates",
    "Two days ago I went back — a first look, not validation",
    "What this can and cannot tell you",
    "The 145 polygons are open data now",
    "Check them on the ground, then put them on the map",
]

# --- 決定性のための固定リテラル値 ---
FIXED_DATETIME = dt.datetime(2026, 1, 1, 0, 0, 0)
FIXED_AUTHOR = "geo-laboratory exp002"
FIXED_LAST_MODIFIED_BY = "exp002_kitagi_foss4g2026_presentation.py"
FIXED_REVISION = 1
FIXED_ZIP_DATETIME = (2026, 1, 1, 0, 0, 0)


def _normalize_zip_timestamps(path: Path) -> None:
    """保存済み .pptx（zip）内の各エントリの mtime を固定値で揃え直す。

    python-pptx（zipfile.ZipFile.writestr）は各エントリのタイムスタンプを
    実行時刻で書き込むため、docProps/core.xml の日時を固定しても
    zip 全体のバイト列は実行するたびに変わってしまう。保存後にこの関数で
    全エントリの date_time を固定し、内容（バイト列・圧縮方式・並び順）は
    保ったまま再書き込みすることで、再生成時の md5 一致（決定性）を保証する。
    """
    with zipfile.ZipFile(path, "r") as zin:
        infos = zin.infolist()
        contents = [zin.read(info.filename) for info in infos]

    buf = io.BytesIO()
    with zipfile.ZipFile(buf, "w", compression=zipfile.ZIP_DEFLATED) as zout:
        for info, data in zip(infos, contents):
            fixed_info = zipfile.ZipInfo(info.filename, date_time=FIXED_ZIP_DATETIME)
            fixed_info.compress_type = info.compress_type
            fixed_info.external_attr = info.external_attr
            zout.writestr(fixed_info, data)

    path.write_bytes(buf.getvalue())


def add_title(slide, text: str):
    """タイトル用テキストボックスを配置する。

    検査（validate_exp002_kitagi_foss4g2026_presentation.py）が図形を確実に
    識別できるよう、シェイプ名を "Title" に固定する。
    """
    box = slide.shapes.add_textbox(
        MARGIN, TITLE_TOP, SLIDE_W - 2 * MARGIN, TITLE_HEIGHT
    )
    box.name = "Title"
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = PP_ALIGN.LEFT
    run = p.runs[0]
    run.font.size = Pt(SZ_TITLE)
    run.font.bold = True
    run.font.color.rgb = COL_TEXT
    return box


def add_body(
    slide,
    lines: list[str | tuple[str, int]],
    size_pt: int = SZ_BODY,
    *,
    top=None,
    left=None,
    width=None,
    height=None,
    align=PP_ALIGN.LEFT,
    color: RGBColor = COL_TEXT,
    bold: bool = False,
    name: str = "Body",
):
    """本文用テキストボックスを配置する（Task 2 で用意、Task 3 で使用）。

    `lines` の各要素は `str`（`size_pt` を使用）または `(text, size_pt)` の2要素
    タプル（行ごとに個別のフォントサイズを指定。evidence 階層で補足値を主結果より
    小さく見せる場合に使う）。いずれの `size_pt` も BODY_MIN_PT（15pt）を
    下回ってはならない。

    位置・寸法（`top`/`left`/`width`/`height`）を省略した場合は Task 2 の既定位置
    （タイトル直下から footer 手前まで）を使う。大きな数値（S6 の `145` 等）を
    60〜72pt の callout として個別のテキストボックスに置く場合は、これらの
    キーワード引数で位置を指定して呼び出す。シェイプ名は既定で "Body"（Task 2 の
    命名規約を維持する）。
    """
    resolved: list[tuple[str, int]] = [
        (line, size_pt) if isinstance(line, str) else line for line in lines
    ]
    for _, pt in resolved:
        if pt < BODY_MIN_PT:
            raise ValueError(f"body size_pt={pt} は下限 {BODY_MIN_PT}pt を下回っている")

    box = slide.shapes.add_textbox(
        left if left is not None else MARGIN,
        top if top is not None else TITLE_TOP + TITLE_HEIGHT,
        width if width is not None else SLIDE_W - 2 * MARGIN,
        height if height is not None else SLIDE_H - TITLE_TOP - TITLE_HEIGHT - Inches(0.5),
    )
    box.name = name
    tf = box.text_frame
    tf.word_wrap = True
    for i, (line, pt) in enumerate(resolved):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.alignment = align
        run = p.runs[0]
        run.font.size = Pt(pt)
        run.font.bold = bold
        run.font.color.rgb = color
    return box


def add_footer(slide, text: str | list[str], size_pt: int = SZ_FOOTER):
    """フッター用テキストボックスを配置する（Task 2 で用意、Task 3 で使用）。

    `text` は単一行の `str`、または複数行を表す `list[str]`（帰属・ライセンス表記と
    ライブラリ名を別行にする場合等）。`size_pt` は 11〜12pt を想定する。
    """
    lines = [text] if isinstance(text, str) else list(text)
    extra = Inches(0.22) * (len(lines) - 1)
    box = slide.shapes.add_textbox(
        MARGIN, SLIDE_H - Inches(0.4) - extra,
        SLIDE_W - 2 * MARGIN - NUM_W - Inches(0.1),
        Inches(0.3) + extra,
    )
    box.name = "Footer"
    tf = box.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.alignment = PP_ALIGN.LEFT
        run = p.runs[0]
        run.font.size = Pt(size_pt)
        run.font.color.rgb = COL_MUTED
    return box


def add_slide_number(slide, n: int):
    """スライド番号を右下に配置する。"""
    box = slide.shapes.add_textbox(
        SLIDE_W - MARGIN - NUM_W, SLIDE_H - Inches(0.4), NUM_W, NUM_H
    )
    box.name = "SlideNumber"
    tf = box.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.text = str(n)
    p.alignment = PP_ALIGN.RIGHT
    run = p.runs[0]
    run.font.size = Pt(SZ_NUM)
    run.font.color.rgb = COL_MUTED
    return box


def s01(slide, n: int) -> None:
    """表紙。演者・カンファレンス・日時（内容契約 Slide 1 Projected body）。

    視覚: 丁場池の写真（`choba_lake_3.jpg`）を右側に大きく配置し、タイトル・
    演者情報は左揃えで残す。

    Fix round 2（レビュー指摘）: 元は `fig03_keirin_cliff.jpg`（日本語記事用に
    グレースケール化された図版）を使っていたが、島の特徴である「緑がかった水」を
    表紙で見せられないのは事故的な欠落であるとのレビュー指摘を受け、色付きの
    `choba_lake_3.jpg`（垂直な花崗岩壁＋緑がかった水面＋青空）に差し替えた
    （担当者ルーリング）。S7 パネル(a)でも同じ写真を再利用する（再認識が目的）。
    """
    add_title(slide, TITLES[0])
    cover_path = IMAGES / "choba_lake_3.jpg"
    photo_h = GRAPHIC_HEIGHT
    cover_iw, cover_ih = _image_size(cover_path)
    photo_w = int(photo_h * cover_iw / cover_ih)  # 元写真のアスペクト比に合わせ、クロップなしで大きく配置
    photo_left = SLIDE_W - MARGIN - photo_w
    add_picture_cover(slide, cover_path, photo_left, GRAPHIC_TOP, photo_w, photo_h, "Picture1")
    text_width = photo_left - MARGIN - IMAGE_TEXT_GAP
    add_body(
        slide,
        [
            "Noboru Otsuka — Geolonia Inc.",
            "FOSS4G 2026 Hiroshima",
            "2026-09-02 13:30 · Himawari",
        ],
        top=GRAPHIC_TOP, width=text_width, height=GRAPHIC_HEIGHT,
    )
    add_slide_number(slide, n)


def s02(slide, n: int) -> None:
    """島と遺産、そして残された池（内容契約 Slide 2 Projected body）。

    視覚: 位置図（`poster_f1_study_area.png`）を右側に配置。

    Final review の Critical 指摘（図中文字の実効サイズ）への対応:
    この図版はポスターからの流用で、パネル(b)の島名ラベルが native 15pt しかない。
    従来は幅6.0inの箱に収めており実効 11.5pt で投影時に読めなかった。
    地理院タイルをネットワーク取得して描く図版であり再生成できないため、
    **配置をほぼ実寸まで拡大**して実効 15pt 以上を満たす（`S02_MAP_SCALE`）。
    その結果テキスト列が約3.8inに狭まるので、本文は 16pt（下限15ptは維持）にする。
    """
    add_title(slide, TITLES[1])
    img_path = IMAGES / "poster_f1_study_area.png"
    img_w_in, img_h_in = picture_native_size_in(img_path)
    img_w = Emu(int(round(914400 * img_w_in * S02_MAP_SCALE)))
    img_h = Emu(int(round(914400 * img_h_in * S02_MAP_SCALE)))
    assert img_h <= GRAPHIC_HEIGHT, "S2: 位置図の高さが図版帯に収まらない"
    img = add_picture_at_scale(
        slide, img_path,
        SLIDE_W - MARGIN - img_w, GRAPHIC_TOP + (GRAPHIC_HEIGHT - img_h) // 2,
        S02_MAP_SCALE, "Picture1",
    )
    text_width = img.left - MARGIN - IMAGE_TEXT_GAP
    add_body(
        slide,
        [
            "Kitagi Island · Kasaoka City, Okayama · Seto Inland Sea",
            "Granite quarried since the early 17th century",
            "127 active quarry sites at the 1957 peak · up to 12,000 residents",
            "Today: two working quarries · about 600–700 residents",
            "Abandoned pits filled with rain and groundwater",
            'National heritage since 2019 — "Stone Islands of Setouchi"',
            "I found no island-wide record of the ponds themselves.",
        ],
        SZ_BODY_NARROW,
        top=GRAPHIC_TOP, width=text_width, height=GRAPHIC_HEIGHT,
    )
    add_slide_number(slide, n)


def s03(slide, n: int) -> None:
    """徒歩スケール：質感（内容契約 Slide 3 Projected body）。

    視覚: 現地写真2枚を同寸で並置。`fig01_lake_stage.jpg` と `choba_lake_1.jpg` を使う。

    Fix round 3（レビュー指摘）: `choba_lake_2.jpg`（丁場の谷を見下ろして遠景の海まで
    写した構図。水面がほぼ写っていない上、クレーンアーム・小屋が写り込む）は投影文
    `Vertical granite walls, cut not weathered` / `Water in an unusual green` を
    支持しないため不採用とし、`choba_lake_1.jpg`（垂直な花崗岩の切削面が緑がかった
    水面へ直接落ち込む構図）に差し替えた。`choba_lake_3.jpg`（表紙・S7パネル(a)で
    採用済み）とは画角が異なるため、表紙の反復にもならない。
    """
    add_title(slide, TITLES[2])
    add_body(
        slide,
        [
            "March 2026 — a drone mapping party on the island",
            "Vertical granite walls, cut not weathered",
            "Water in an unusual green; reported depths of a few metres to about twenty",
            "A wooden stage on the water, with stone blocks arranged as seats",
            "Five or six sites during the event.",
        ],
        top=GRAPHIC_TOP, height=PHOTO_TEXT_HEIGHT,
    )
    # fig01: 縦長写真をそのまま中央対称クロップすると、水面が写る下部が失われ岩壁のみに
    # なる（元画像は上2/3が岩壁・下1/3が水面のため）。水面・ステージが見える構図を
    # 残すよう上側を多くクロップする（vbias1=0.75、Task 4 から変更なし）。
    # choba_lake_1: 上部は空・稜線の樹木、中央から下は垂直な花崗岩壁、最下部が緑がかった
    # 水面（元画像1500x1999で水面は下から約22%）。中央対称クロップでは空・樹木が残って
    # しまうため、上側を大きくクロップして水面を残す（vbias2=0.85。0.70〜1.00で試作
    # クロップを目視確認し、空・稜線の崩石が完全に外れ、垂直な壁面と水面がともに
    # 十分な面積で残る構図として選定）。
    add_photo_pair(
        slide, IMAGES / "fig01_lake_stage.jpg", IMAGES / "choba_lake_1.jpg",
        vbias1=0.75, vbias2=0.85,
    )
    add_slide_number(slide, n)


def s04(slide, n: int) -> None:
    """上空スケール：境界（内容契約 Slide 4 Projected body）。

    視覚: 丁場池の上空写真（`aerial_quarry_pond.jpg`）と、湖上ステージのドローン
    （`drone_lake_stage.jpg`）を**縦長スロット2枚**として同寸で並置し、本文列を左に置く。

    2026-08-24 の調整: 発表者から色付き原本の提供を受け、印刷用にグレースケール化されて
    いた記事図版（`fig06_aerial_quarries.jpg` / `fig05_drone_takeoff.jpg`）を置き換えた。
    新しい2枚はいずれも縦位置（縦横比 0.75〜0.88）で、従来の 16:9 スロットでは各フレームの
    58% を捨てることになる。S4 の主張（採石権の境界が地形として読める）は縦方向の広がりを
    要するため、スロットを縦長にし、本文を上ではなく左に置くレイアウトへ変更した
    （S3・S9 の 16:9 固定スロット `PHOTO_SLOT_*` は変更していない）。
    実測: 縦長スロット 3.87 x 4.40 in を2枚並べると写真帯は 7.93 in、本文列は 4.00 in と
    なり、図版帯（12.23 x 4.80 in）に重なりなく収まる。
    """
    add_title(slide, TITLES[3])
    slot_w = Emu(int(round(S04_PHOTO_SLOT_H * S04_PHOTO_SLOT_AR)))
    photos_w = 2 * slot_w + S04_PHOTO_GAP
    photos_left = SLIDE_W - MARGIN - photos_w
    photos_top = GRAPHIC_TOP + (GRAPHIC_HEIGHT - S04_PHOTO_SLOT_H) // 2
    assert photos_left >= USABLE_LEFT, "S4: 縦長スロット2枚が本文列の領域に食い込んでいる"
    assert photos_top >= GRAPHIC_TOP and photos_top + S04_PHOTO_SLOT_H <= GRAPHIC_BOTTOM, (
        "S4: 縦長スロットが図版帯の上下に収まらない"
    )
    text_width = photos_left - MARGIN - IMAGE_TEXT_GAP
    add_body(
        slide,
        [
            "Drone flown from the stage on the water",
            "Grey rectangles cut into the green canopy",
            "Between two quarries, a thin wall left standing",
            "A property line, standing as terrain",
            "The same event added features to OpenStreetMap.",
        ],
        top=GRAPHIC_TOP, width=text_width, height=GRAPHIC_HEIGHT,
    )
    # 左: 丁場池の上空写真。スロットの縦横比を実比に合わせてあるためクロップは生じない。
    add_picture_cover(
        slide, IMAGES / "aerial_quarry_pond.jpg",
        photos_left, photos_top, slot_w, S04_PHOTO_SLOT_H, "Picture1",
    )
    # 右: 湖上ステージのドローン。1080 x 1440 px はスロットより縦長なので上下クロップに
    # なる。被写体（岩壁・水面・機体2機）は上 65% に収まっており、下部は被写体の無い
    # 板張りの床なので、クロップは全量を下側から取る（vbias=0.0）。
    add_picture_cover(
        slide, IMAGES / "drone_lake_stage.jpg",
        photos_left + slot_w + S04_PHOTO_GAP, photos_top, slot_w, S04_PHOTO_SLOT_H,
        "Picture2", vbias=S04_PHOTO_VBIAS_2,
    )
    add_slide_number(slide, n)


def s05(slide, n: int) -> None:
    """衛星スケール：手法（内容契約 Slide 5 Projected body）。

    ルーリング（NDWI・MNDWI・NDVI の式定義そのものはノートと図版内へ）に従い、
    投影本文には合成条件（`Water if ...`）のみを置き、`(Green − NIR)` 等の個別式は
    ここに書かない（図版 P5 のパネル内注記が指数名と閾値を示す）。

    Final review の Critical 指摘への対応: 従来置いていた `poster_f4_index_panels.png`
    は native 18pt・実寸 8.82 x 8.50 in で、16:9スライドのどんな配置でも実効 15pt に
    届かない（15pt には配置高さ 7.09in が必要で、スライド全高は 7.5in）。
    そのため、同じパネル画像を切り出して英語ラベルを大きく描き直した
    `p05_index_panels.png`（`exp002_kitagi_foss4g2026_figures.py` の
    `make_p05_index_panels()` が生成。ネットワーク非依存）に差し替えた。
    """
    add_title(slide, TITLES[4])
    img_max_w = Inches(5.4)
    img = add_picture_contain(
        slide, IMAGES / "p05_index_panels.png",
        SLIDE_W - MARGIN - img_max_w, GRAPHIC_TOP, img_max_w, GRAPHIC_HEIGHT, "Picture1",
    )
    text_width = img.left - MARGIN - IMAGE_TEXT_GAP
    add_body(
        slide,
        [
            "Sentinel-2 L2A via the Microsoft Planetary Computer STAC API",
            "Summer 2025-08-02 · 0.7% cloud · 10 m analysis grid",
            "Water if (NDWI > −0.2 OR MNDWI > −0.1) AND NOT (NDVI > 0.3)",
            "Thresholds below zero: at 10 m a narrow pond is part water, part granite, part shadow",
            "A standard water-index workflow — nothing new in the method",
            "Minimum reported polygon area: 100 m²",
        ],
        top=GRAPHIC_TOP, width=text_width, height=GRAPHIC_HEIGHT,
    )
    add_slide_number(slide, n)


def s06(slide, n: int) -> None:
    """主結果：145件（内容契約 Slide 6 Projected body）。

    evidence 階層: 主結果 `145` を callout（60〜72pt、ここでは SZ_CALLOUT=66pt）、
    春季の補足値 `113` は小タイルとして本文の下限 15pt に留め、主結果より明確に
    小さくする（validator の `check_evidence_hierarchy` が検査する）。
    """
    add_title(slide, TITLES[5])
    img_max_w = Inches(5.2)
    # 図版帯の高さいっぱいに置く（4.4in → 4.8in）。P6 は縦横比がほぼ1で高さ拘束のため、
    # 実効ptを稼げるのは高さ方向だけである（Final review の Critical 指摘への対応）。
    img_height = GRAPHIC_HEIGHT
    img = add_picture_contain(
        slide, IMAGES / "p06_clusters_map.png",
        SLIDE_W - MARGIN - img_max_w, GRAPHIC_TOP, img_max_w, img_height, "Picture1",
    )
    text_width = img.left - MARGIN - IMAGE_TEXT_GAP
    add_body(
        slide, ["145"], SZ_CALLOUT,
        top=TITLE_TOP + TITLE_HEIGHT, left=MARGIN,
        width=Inches(3.2), height=Inches(1.3),
        bold=True, color=COL_ACCENT,
    )
    add_body(
        slide,
        [
            "intra-island water polygons ≥ 100 m², summer 2025-08-02",
            "Clustered in the north, south-east, centre and west — consistent with historical quarrying records",
            "145 detections vs 127 recorded quarry sites — a comparison of scale, not a one-to-one match",
            ("Spring 2025-03-23 — 113 polygons reported", BODY_MIN_PT),
            "These are detected water polygons, not individually field-confirmed quarry ponds.",
        ],
        top=TITLE_TOP + TITLE_HEIGHT + Inches(1.4),
        width=text_width,
        height=Inches(3.0),
    )
    add_footer(slide, "Contains modified Copernicus Sentinel data [2025].")
    add_slide_number(slide, n)


def s07(slide, n: int) -> None:
    """三つの縮尺が示すもの（内容契約 Slide 7 Projected body）。

    視覚: 三スケール合成図（`p07_three_scales.png`、英語ラベルのみ）を本文の下に
    横断的に配置する。

    Task 4 レビュー指摘の修正: 元は `fig09_multiscale.png`（日本語記事
    `docs/articles/2026_chiri-koryu-10/` と共有される図版で、パネル注記が日本語で
    焼き込まれ、(b) パネルに動画UIの写り込みがある）を使っていたが、英語のみの
    投影面という契約に抵触するため、発表専用の英語版図版に差し替えた
    （`exp002_kitagi_foss4g2026_figures.py` の `make_p07_three_scales()` が生成。
    日本語記事側の図版・生成スクリプトは変更していない）。
    """
    add_title(slide, TITLES[6])
    text_h = Inches(1.3)
    add_body(
        slide,
        [
            "On foot — texture: the cut face, the water, the depth",
            "From the air — boundaries: property lines standing as rock walls",
            "From orbit — distribution: 145 candidates across the island",
            "Not better or worse. Different things become visible.",
        ],
        top=GRAPHIC_TOP, height=text_h,
    )
    img_top = GRAPHIC_TOP + text_h + Inches(0.15)
    img_height = GRAPHIC_BOTTOM - img_top
    add_picture_contain(
        slide, IMAGES / "p07_three_scales.png",
        USABLE_LEFT, img_top, USABLE_WIDTH, img_height, "Picture1",
    )
    add_slide_number(slide, n)


def s08(slide, n: int) -> None:
    """5〜6か所 対 145件：規模の対比（内容契約 Slide 8 Projected body）。

    「5〜6か所訪問」対「145件検出」という規模の対比を、2つの callout（60〜72pt）
    として並置する。
    """
    add_title(slide, TITLES[7])
    img_max_w = Inches(4.6)
    # 図版帯の高さいっぱいに置く（4.3in → 4.8in）。理由は S6 と同じ（高さ拘束）。
    img_height = GRAPHIC_HEIGHT
    img = add_picture_contain(
        slide, IMAGES / "p08_visit_anchors_map.png",
        SLIDE_W - MARGIN - img_max_w, GRAPHIC_TOP, img_max_w, img_height, "Picture1",
    )
    text_width = img.left - MARGIN - IMAGE_TEXT_GAP
    add_body(
        slide, ["5–6"], SZ_CALLOUT,
        top=TITLE_TOP + TITLE_HEIGHT, left=MARGIN,
        width=Inches(2.2), height=Inches(1.1),
        bold=True, color=COL_ACCENT,
    )
    add_body(
        slide, ["145"], SZ_CALLOUT,
        top=TITLE_TOP + TITLE_HEIGHT, left=MARGIN + Inches(2.4),
        width=Inches(2.2), height=Inches(1.1),
        bold=True, color=COL_ACCENT,
    )
    add_body(
        slide,
        [
            "Five or six quarry sites visited during the event",
            "145 water polygons detected from one scene",
            "Individual ponds are not field-confirmed — no precision or recall yet",
            "Every quarry feature already mapped in OpenStreetMap overlaps one of the detections (retrieved 2026-08-23, for reference)",
            "The candidates form a finite field-check list.",
        ],
        top=TITLE_TOP + TITLE_HEIGHT + Inches(1.3),
        width=text_width,
        height=Inches(3.0),
    )
    add_slide_number(slide, n)


def s09(slide, n: int) -> None:
    """再訪（内容契約 Slide 9 Projected body）。

    視覚: 8/31 撮影分の写真2枚（撮影前は `resolve_revisit_photo()` がプレースホルダを
    返す）。`--no-revisit` 指定時は `build()` がこのスライド自体を除外する。
    """
    add_title(slide, TITLES[8])
    add_body(
        slide,
        [
            "2026-08-31 — return visit",
            "Candidates selected from the published GeoJSON",
            "Illustrative field photographs — not accuracy validation",
            "What the scan pointed at, seen from the ground.",
        ],
        top=GRAPHIC_TOP, height=PHOTO_TEXT_HEIGHT,
    )
    add_photo_pair(slide, resolve_revisit_photo(1), resolve_revisit_photo(2))
    add_slide_number(slide, n)


def s10(slide, n: int) -> None:
    """限界と次の一手（内容契約 Slide 10 Projected body）。"""
    add_title(slide, TITLES[9])
    add_body(slide, [
        "10 m resolution — ponds narrower than about 10 m are unreliable",
        "Thresholds below zero admit dark rock and shadow",
        "No precision or recall — field validation not done",
        "Next: walk the candidates · a land mask for the shoreline · higher-resolution imagery",
    ])
    add_slide_number(slide, n)


def s11(slide, n: int) -> None:
    """オープンデータ（内容契約 Slide 11 Projected body）。

    帰属・ライセンス・使用ライブラリ名は `add_footer` へ送る（11〜12pt）。
    """
    add_title(slide, TITLES[10])
    qr_size = Inches(2.6)
    img = add_picture_contain(
        slide, IMAGES / "poster_qr_repo.png",
        SLIDE_W - MARGIN - qr_size, GRAPHIC_TOP, qr_size, qr_size, "Picture1",
    )
    text_width = img.left - MARGIN - IMAGE_TEXT_GAP
    add_body(
        slide,
        [
            "Published — 145 detected polygons as GeoJSON · EPSG:4326",
            "Pipeline outputs — GeoJSON and GeoTIFF for fieldwork and heritage documentation",
            "Open-source Python pipeline · no licence fee, no imagery purchase",
            "The same workflow could be extended to other quarried islands in the Seto Inland Sea.",
        ],
        width=text_width,
        height=Inches(4.4),
    )
    add_footer(slide, [
        "rasterio · numpy · shapely · pystac-client · planetary-computer · folium",
        "Contains modified Copernicus Sentinel data [2025]. CC BY 4.0.",
        # 基図の帰属はデッキが実際に使っている2種類だけを挙げる（他は挙げない）。
        #   S2 の位置図（`poster_f1_study_area.png`）= 地理院タイル英語版
        #   S6・S7 パネル(c)・S8 の検出地図 = CARTO Positron（OpenStreetMap データ）
        # 1行に収めると 11pt でも折り返して行がスライド下端へ迫るため、独立した行にする。
        "Basemaps: GSI Tiles, Geospatial Information Authority of Japan; "
        "© OpenStreetMap contributors, © CARTO.",
    ])
    add_slide_number(slide, n)


def s12(slide, n: int) -> None:
    """地図への還元（内容契約 Slide 12 Projected body）。

    視覚: 「衛星 → 現地 → 地図」の3ステップフロー（`p12_loop_diagram.png`）を
    本文の下に横断的に配置。
    """
    add_title(slide, TITLES[11])
    text_h = Inches(1.9)
    add_body(
        slide,
        [
            "Satellite scan → a finite candidate list",
            "Field visit → see it with your own eyes",
            "OpenStreetMap → put what you confirmed on the public map",
            "I plan to contribute the ponds I can confirm.",
            "The March mapping party added features observed on the ground. The scan suggests where to look next.",
            "Thank you · Q&A",
        ],
        top=GRAPHIC_TOP, height=text_h,
    )
    img_top = GRAPHIC_TOP + text_h + Inches(0.15)
    img_height = GRAPHIC_BOTTOM - img_top
    add_picture_contain(
        slide, IMAGES / "p12_loop_diagram.png",
        USABLE_LEFT, img_top, USABLE_WIDTH, img_height, "Picture1",
    )
    add_slide_number(slide, n)


def _strip_section_rule(text: str) -> str:
    """節末の水平線（Markdown の `---`）を落とす。

    節の区切りは Markdown の可読性のためのもので、ノートペインに載せる内容ではない。
    """
    return re.sub(r"\n+-{3,}\s*$", "", text.strip()).strip()


def parse_speaker_notes() -> dict[int, dict[str, str]]:
    """スピーカーノート Markdown を内容契約のスライド番号ごとに読み取る。

    `### Slide N — <title>` 見出しで節に分け、`**EN (spoken)**` から
    `**JA (not spoken)**` の直前までを英語の発話本文、`**JA (not spoken)**` の行から
    節末までを日本語の非発話部分として取り出す。
    """
    md = NOTES_MD.read_text(encoding="utf-8")
    parts = re.split(r"(?m)^### Slide (\d+) — (.*)$", md)[1:]
    sections: dict[int, dict[str, str]] = {}
    for i in range(0, len(parts), 3):
        number = int(parts[i])
        body = parts[i + 2]
        if EN_MARKER not in body or JA_MARKER not in body:
            raise ValueError(f"ノートの Slide {number} に EN / JA の区切りが無い")
        after_en = body.split(EN_MARKER, 1)[1]
        en = after_en.split(JA_MARKER, 1)[0].strip()
        ja = _strip_section_rule(JA_MARKER + after_en.split(JA_MARKER, 1)[1])
        sections[number] = {"title": parts[i + 1].strip(), "en": en, "ja": ja}
    return sections


def build_notes_text(section: dict[str, str]) -> str:
    """ノートペインに書き込む文字列を組み立てる（英語→空行→日本語）。

    日本語側は見出し行（`**JA (not spoken)**` 以下）を含めたまま入れる。発表者が
    ノートペインだけを見た状態でも「ここから先は読み上げない」境界が分かるようにする。
    """
    return f"{section['en']}\n\n{section['ja']}"


def build(revisit: bool = True) -> Presentation:
    """12枚（`revisit=False` の場合は S9 を除いた11枚）の Presentation を組み立てる。

    `revisit` は CLI の `--no-revisit`（`main()` 参照）から渡される。

    スライド番号（投影面の右下）はデッキ内の連番だが、スピーカーノートは内容契約の
    スライド番号で引く。再訪なし版では S9 が抜けて両者がずれるため、
    `(契約番号, 生成関数)` の組で持ち、番号の用途を分離する。
    """
    prs = Presentation()
    prs.slide_width, prs.slide_height = SLIDE_W, SLIDE_H

    builders: list[tuple[int, object]] = [
        (1, s01), (2, s02), (3, s03), (4, s04), (5, s05), (6, s06), (7, s07), (8, s08),
    ]
    if revisit:
        builders.append((9, s09))
    builders += [(10, s10), (11, s11), (12, s12)]

    notes = parse_speaker_notes()

    for position, (contract_n, fn) in enumerate(builders, start=1):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        fn(slide, position)
        if contract_n not in notes:
            raise ValueError(f"ノートに Slide {contract_n} の節が無い")
        slide.notes_slide.notes_text_frame.text = build_notes_text(notes[contract_n])

    cp = prs.core_properties
    cp.created = FIXED_DATETIME
    cp.modified = FIXED_DATETIME
    cp.last_modified_by = FIXED_LAST_MODIFIED_BY
    cp.revision = FIXED_REVISION
    cp.author = FIXED_AUTHOR
    cp.title = TITLES[0]

    return prs


def main() -> None:
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument(
        "--no-revisit",
        action="store_true",
        help=(
            "S9（8/31 再訪スライド）を除いた11枚を "
            "exp002_kitagi_foss4g2026_presentation_no_revisit.pptx として出力する"
        ),
    )
    args = parser.parse_args()

    revisit = not args.no_revisit
    out = NO_REVISIT_OUTPUT if args.no_revisit else OUTPUT
    prs = build(revisit=revisit)
    prs.save(out)
    _normalize_zip_timestamps(out)
    print(f"saved: {out} ({len(prs.slides)} slides)")


if __name__ == "__main__":
    main()
