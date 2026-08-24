#!/usr/bin/env python3
"""FOSS4G 2026 北木島 口頭発表の成果物を検査する。

投影面に関する検査群（スライド数・スライド寸法・タイトル文字列と28pt・必須文字列・
禁止表現・CJK混入・evidence階層・本文15pt下限・footer 11〜12pt・画像数・callout の
文字サイズ・pin した画像のバイト一致・**図中文字の実効サイズ**・実効解像度・
装飾効果の不使用）は、**12枚版と再訪なし版（11枚）の両方**に対して同じものを走らせる
（`main()` の `DECKS`）。以前は再訪なし版がスライド数と1文字列しか検査されていなかった。

図中文字の実効サイズは、生成済みPPTXの `shape.width` と画像実寸（px ÷ dpi）から
配置倍率を復元し、図版生成スクリプトが宣言した native サイズ（`NATIVE_FONT_SIZES`）に
掛けて求める（`check_placed_font_sizes()`）。配置幅を仮定した自己申告は使わない。

このほか、スピーカーノート Markdown（構造・英語語数・S6 の required spoken content・
再訪ノートの申し送り）と PPTX ノートペインとの同期、Task 6 の照合記録
（`exp002_kitagi_foss4g2026_presentation_verification.md`）の数値の出典・
images/ 全ファイルのSHA256・8/31撮影テーブルの見出しを検査する。

使い方: uv run python docs/presentations/validate_exp002_kitagi_foss4g2026_presentation.py
"""

from __future__ import annotations

import hashlib
import re
import sys
import zipfile
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

# 図版の native フォントサイズ宣言と 15pt 下限は図版生成スクリプトを正本とする
# （検査側に数表を複製すると、複製の方が古くなっても検査が通ってしまう）。
from exp002_kitagi_foss4g2026_figures import NATIVE_FONT_SIZES, SLIDE_PT_FLOOR

BASE = Path(__file__).resolve().parent
IMAGES = BASE / "images"
PPTX = BASE / "exp002_kitagi_foss4g2026_presentation.pptx"
NO_REVISIT_PPTX = BASE / "exp002_kitagi_foss4g2026_presentation_no_revisit.pptx"
NOTES_MD = BASE / "exp002_kitagi_foss4g2026_presentation_speaker_notes.md"
VERIFICATION = BASE / "exp002_kitagi_foss4g2026_presentation_verification.md"

# スピーカーノートの英語発話量の判定基準。内容契約「タイミング」節の秒数（合計1,050秒）を
# 145 wpm で語数へ換算し、±25% の幅で判定する。
WPM = 145
DURATIONS_S = [35, 100, 90, 90, 110, 150, 70, 90, 110, 80, 70, 55]
WORD_TOLERANCE = 0.25

# ノート Markdown 内の構造マーカー。英語（発話対象）と日本語（非発話）の境界。
EN_MARKER = "**EN (spoken)**"
JA_MARKER = "**JA (not spoken)**"

# 内容契約 Slide 6 の `Required spoken content`。内部注記ではなく英語の発話本文に
# 置くことが契約上の要件であるため、S6 の EN 部分に対して検査する。
#
# 契約が要求しているのは**事実が発話されること**であり、スライドのタイル文字列を
# そのまま読み上げることではない。逐語の部分文字列で固定すると
# `Spring: 2025-03-23, 0.0% cloud — 113 reported polygons` のような
# 動詞のない・日付が読み上げられない書き言葉を強制してしまうため、事実単位の
# 正規表現で検査し、数字表記と綴り字表記の双方を受け入れる。
# (事実の説明, その事実が発話されていると認める正規表現) の並び。
S6_REQUIRED_FACTS: list[tuple[str, str]] = [
    ("春季シーンの日付", r"2025-03-23|twenty-third of March"),
    # 雲量0.0%は**数値として発話される**ことを要求する。`no cloud at all` のような
    # 見た目の印象に置き換えると「報告されたメタデータの値」であることが伝わらないため、
    # 数値表記と綴り字表記（zero point zero percent）のみを受け入れる。
    ("春季シーンの雲量（0.0%）", r"0\.0%|zero point zero percent"),
    ("春季の報告件数（113）", r"\b113\b|one hundred and thirteen"),
    ("春季の最大面積（1.28 ha）", r"1\.28 hectares|one point two eight hectares"),
    ("夏季シーンの日付", r"2025-08-02|second of August"),
    ("夏季シーンの雲量（0.7%）", r"0\.7%|zero point seven percent"),
    ("夏季の検出件数（145）", r"\b145\b|one hundred and forty-five"),
    ("夏季の最大面積（7,826 m²）",
     r"7,826 square metres|seven thousand eight hundred and twenty-six square metres"),
    ("差の原因が特定できていないこと", r"have not isolated the cause"),
    ("春季実行の設定が保存されていないこと", r"configuration is not preserved"),
    ("NDVIマスクの除外が9ピクセルであること", r"removed only nine pixels"),
    ("候補であって個別の現地確認をしていないこと",
     r"not individually field-confirmed quarry ponds"),
]

# S9 の再訪写真スロットに実写が入ったかどうかの判定材料。生成スクリプトの
# `resolve_revisit_photo()` と同じファイル名・同じ拡張子の並びを見る。
REVISIT_PHOTO_STEMS = ("revisit_1", "revisit_2")
REVISIT_PHOTO_EXTS = ("jpg", "jpeg", "png")

# 再訪写真がプレースホルダのままである間、S9 のノートに残すことを必須とする目印。
# S9 の英語本文は再訪前に書いたものであり、写真の差し替え時に本文も見直す必要がある。
REVISIT_UPDATE_MARKER = "[UPDATE AFTER 2026-08-31]"

# 再訪なし版（11枚）の各スライド位置に対応する内容契約のスライド番号。
# S9（再訪）が抜けるため、位置9以降は内容契約の番号とずれる。
NO_REVISIT_CONTRACT_NUMBERS = [1, 2, 3, 4, 5, 6, 7, 8, 10, 11, 12]

# 投影時の実効解像度の下限（dpi）。「配置幅あたりのピクセル数」で判定する。
# 以前はここに「配置幅220mmで200dpiを満たす最小ピクセル幅（1732px）」という定数が
# あったが、220mm はデッキが実際には使っていない配置幅であり、図中フォントの
# 自己申告と同じ誤りだった（Final review の Critical 指摘）。実配置は
# 生成済みPPTXの `shape.width` から読む（`check_placed_resolution`）。
MIN_PLACED_DPI = 200.0

FIGURES = (
    "p05_index_panels.png", "p06_clusters_map.png", "p07_three_scales.png",
    "p08_visit_anchors_map.png", "p12_loop_diagram.png",
)

# バイト単位で固定する写真・図版のソース。キーはスライド番号、値はそのスライドの
# PICTURE シェイプの挿入順で期待するファイル名（`docs/presentations/images/`）。
# S1・S3 は色付きの現地写真を意図しており、日本語記事用にグレースケール化された
# `fig03_keirin_cliff.jpg` へ将来の編集で意図せず戻ってしまう回帰を防ぐ
# （Fix round 2 のレビュー指摘）。S7 は英語ラベルのみの三スケール合成図
# `p07_three_scales.png` が、日本語キャプション付きの `fig09_multiscale.png` に
# 戻ってしまう回帰を防ぐ（Fix round 1 のレビュー指摘）。
# S4 は 2026-08-24 の調整で、印刷用にグレースケール化された記事図版
# （`fig06_aerial_quarries.jpg` / `fig05_drone_takeoff.jpg`、いずれも削除済み）から
# 発表者提供の色付き原本へ差し替えた。グレースケール版へのフォールバックも、
# macOS の Dock / ツールチップを含む未クロップの原本の混入も、バイト比較で検出する。
PINNED_PHOTO_SOURCES: dict[int, tuple[str, ...]] = {
    1: ("choba_lake_3.jpg",),
    3: ("fig01_lake_stage.jpg", "choba_lake_1.jpg"),
    4: ("aerial_quarry_pond.jpg", "drone_lake_stage.jpg"),
    7: ("p07_three_scales.png",),
}

# 各スライドの画像（PICTURE シェイプ）の期待枚数。S9 は8/31撮影分（未着の間はプレースホルダ）
# 2枚、S10 はテキスト中心のため意図的に0枚。
EXPECTED_IMAGES: dict[int, int] = {
    1: 1, 2: 1, 3: 2, 4: 2, 5: 1, 6: 1, 7: 1, 8: 1, 9: 2, 11: 1, 12: 1,
}

# callout（S6 の `145`、S8 の `5–6`・`145`）が満たすべき文字サイズの範囲（pt）。
CALLOUT_MIN_PT = 60.0
CALLOUT_MAX_PT = 72.0
CALLOUT_TARGETS: dict[int, list[str]] = {6: ["145"], 8: ["5–6", "145"]}

# 内容契約（docs/presentations/exp002_kitagi_foss4g2026_presentation.md）の
# `## Slide N — ` 見出しから転記した、各スライドのタイトル完全一致文字列。
TITLES = [
    "Detecting Quarry Pond Remnants on a Japanese Island Heritage Site Using Sentinel-2 Imagery and Open-Source Remote Sensing Tools",
    "A quarrying island, and the ponds it left behind",
    "On foot: I could stand in front of five or six of them",
    "From the air: the quarry boundaries are the landform",
    "On the train home: one satellite scene covers the whole island",
    "The scan found 145 water polygons across the island",
    "Each scale shows what the others cannot",
    "Five or six sites visited — the scan produced 145 candidates",
    "Two days ago I went back — a first look, not validation",
    "What this can and cannot tell you",
    "The 145 polygons are open data now",
    "Check them on the ground, then put them on the map",
]

# 内容契約の `Projected body` から転記した、スライド別の必須文字列（機械検査対象の部分集合）。
REQUIRED_STRINGS: dict[int, list[str]] = {
    1: ["Noboru Otsuka", "Geolonia", "FOSS4G 2026 Hiroshima", "2026-09-02 13:30", "Himawari"],
    2: ["127 active quarry sites at the 1957 peak", "National heritage since 2019",
        "I found no island-wide record of the ponds themselves"],
    3: ["Five or six sites during the event"],
    4: ["A property line, standing as terrain"],
    5: ["Microsoft Planetary Computer STAC API", "NDWI", "MNDWI", "NDVI",
        "A standard water-index workflow", "100 m"],
    6: ["145", "2025-08-02", "113", "a comparison of scale, not a one-to-one match",
        "not individually field-confirmed quarry ponds",
        "Contains modified Copernicus Sentinel data [2025]"],
    7: ["Not better or worse"],
    8: ["Five or six quarry sites visited during the event",
        "145 water polygons detected from one scene",
        "no precision or recall yet",
        "Every quarry feature already mapped in OpenStreetMap overlaps one of the detections",
        "for reference", "finite field-check list"],
    9: ["2026-08-31", "Illustrative field photographs — not accuracy validation"],
    10: ["10 m resolution", "No precision or recall", "field validation not done"],
    11: ["Published — 145 detected polygons as GeoJSON", "Pipeline outputs",
         "Open-source Python pipeline", "Seto Inland Sea", "CC BY 4.0"],
    12: ["Satellite scan", "Field visit", "OpenStreetMap",
         "I plan to contribute the ponds I can confirm",
         "The March mapping party added features observed on the ground"],
}

# (禁止パターン, 許容される否定形パターン（None なら常に禁止）, 説明)。
# 否定形にマッチする範囲に完全に包含される禁止パターンの一致は許容する
# （例: "not individually field-confirmed quarry ponds" は許容し、
#  裸の "confirmed quarry ponds" は禁止する）。
FORBIDDEN: list[tuple[str, str | None, str]] = [
    (r"confirmed quarry pond", r"not\s+(?:individually\s+)?(?:field[-\s])?confirmed quarry pond",
     "現地確認済みと読める表現"),
    (r"\b145 quarry ponds?\b", None, "検出を丁場池と同一視"),
    (r"high accuracy|accuracy is high", None, "精度主張"),
    (r"\bposter\b", None, "採択形式は口頭発表"),
    (r"one-to-one match", r"not\s+(?:a\s+)?one-to-one match", "1対1対応と読める表現"),
]

BODY_MIN_PT = 15  # 生成スクリプト側の下限（exp002_kitagi_foss4g2026_presentation.py と同値）

# 16:9 スライドの EMU 寸法（生成スクリプトの SLIDE_W / SLIDE_H と同値）。
SLIDE_W_EMU = 12192000
SLIDE_H_EMU = 6858000

# タイトルと footer / スライド番号の文字サイズ（生成スクリプトの SZ_TITLE・SZ_FOOTER・SZ_NUM）。
TITLE_PT = 28.0
FOOTER_PT_RANGE = (11.0, 12.0)

# 投影面は英語のみ（内容契約）。日本語が焼き込まれた図版・訳文の混入を検出する。
# `·` `–` `—` `−` `≥` `²` などの非ASCII約物は既に投影面で使われているため、
# 「非ASCII禁止」ではなく CJK ブロック（漢字・かな・全角約物・全角英数）を禁止する。
CJK_PATTERN = re.compile(
    "["
    "\u3000-\u303f"  # CJK 約物（、。「」等）
    "\u3040-\u309f"  # ひらがな
    "\u30a0-\u30ff"  # カタカナ
    "\u4e00-\u9fff"  # CJK 統合漢字
    "\uff00-\uffef"  # 全角英数・半角カナ
    "]"
)

# 装飾効果の禁止（DESIGN_GUIDE: 影・グラデーションを使わない）。
# 検査対象は作図した内容そのもの（`ppt/slides/*.xml`）に限る。python-pptx の既定
# テーマ `ppt/theme/theme1.xml` は fmtScheme に gradFill / outerShdw を含むため、
# パッケージ全体を対象にすると常に失敗してしまう。
FORBIDDEN_XML_EFFECTS = ("outerShdw", "gradFill")

errors: list[str] = []
checks = 0


def check(cond: bool, msg: str) -> None:
    global checks
    checks += 1
    if not cond:
        errors.append(msg)


def check_figures() -> None:
    """自前生成の図版が存在し、実寸計算に必要な dpi が記録されていることを確認する。

    投影時の精細さ（実効dpi）は配置幅に依存するため、ここではなく
    `check_placed_resolution()` が生成済みPPTXの配置幅から判定する。
    """
    for name in FIGURES:
        path = IMAGES / name
        check(path.is_file(), f"図版が存在しない: {name}")
        if not path.is_file():
            continue
        with Image.open(path) as img:
            width, height = img.size
            dpi = img.info.get("dpi")
        check(width > 0 and height > 0, f"{name}: 画像寸法が不正")
        check(
            dpi is not None and dpi[0] > 0,
            f"{name}: PNG に dpi が記録されていない（実寸が計算できない）",
        )


def deck_slides(label: str, path: Path, contract_numbers: list[int]) -> dict[int, object] | None:
    """デッキを開き、内容契約のスライド番号 → スライドの対応表を返す。

    12枚版は「位置 = 契約番号」だが、再訪なし版（11枚）は S9 が抜けるため位置が
    ずれる。以降の検査はすべてこの対応表を通して契約番号で参照する
    （位置で素朴に添字を取ると、S10 の位置に S9 の内容が載っていても気づけない）。
    スライド数とスライド寸法もここで確認する。
    """
    check(path.is_file(), f"{label}: PPTX が存在しない")
    if not path.is_file():
        return None
    prs = Presentation(path)
    slides = list(prs.slides)
    check(
        len(slides) == len(contract_numbers),
        f"{label}: スライド数が {len(contract_numbers)} でない: {len(slides)}",
    )
    check(
        prs.slide_width == SLIDE_W_EMU and prs.slide_height == SLIDE_H_EMU,
        f"{label}: スライド寸法が {SLIDE_W_EMU}x{SLIDE_H_EMU} EMU (16:9) でない: "
        f"{prs.slide_width}x{prs.slide_height}",
    )
    return {cn: slides[i] for i, cn in enumerate(contract_numbers) if i < len(slides)}


def check_titles(label: str, by_contract: dict[int, object]) -> None:
    """各スライドのタイトル用シェイプ（name="Title"）が内容契約のタイトル文字列と
    完全一致し、数字で始まっておらず、28pt であることを確認する。

    ブリーフの原案（`title in texts[i]` による部分一致＋番号プレフィックスの雑な検出）は
    誤検出・見逃しの両方向で信頼できないため採用しない。代わりにタイトル用シェイプを
    name="Title" で明示的に識別し、テキストの完全一致と先頭文字が数字でないことを
    それぞれ独立に検査する。
    """
    for n, expected in enumerate(TITLES, start=1):
        slide = by_contract.get(n)
        if slide is None:
            continue  # このデッキに含まれないスライド（再訪なし版の S9）
        tag = f"{label} S{n}"
        title_shapes = [
            sh for sh in slide.shapes
            if sh.has_text_frame and sh.name == "Title"
        ]
        check(
            len(title_shapes) == 1,
            f"{tag}: タイトル用シェイプ(name='Title')が1つでない: {len(title_shapes)}個",
        )
        if len(title_shapes) != 1:
            continue
        text = title_shapes[0].text_frame.text
        check(text == expected, f"{tag}: タイトルが完全一致しない — 期待 '{expected[:40]}...'")
        check(not text[:1].isdigit(), f"{tag}: タイトル先頭が数字で始まっている")
        sizes = [
            run.font.size.pt
            for para in title_shapes[0].text_frame.paragraphs
            for run in para.runs
            if run.font.size is not None
        ]
        check(bool(sizes), f"{tag}: タイトルの文字サイズが設定されていない")
        for pt in sizes:
            check(pt == TITLE_PT, f"{tag}: タイトルが {pt}pt — {TITLE_PT:g}pt でない")


def check_footer_sizes(label: str, by_contract: dict[int, object]) -> None:
    """footer とスライド番号のランが 11〜12pt の範囲にあることを確認する。

    本文下限（15pt）の例外として `check_font_floor` が 12pt 以下を許容しているのは
    「footer とスライド番号として意図されたもの」だけである。その意図が本当に
    守られているか（例: 本文が誤って footer 用シェイプへ流れ込んでいないか）を、
    シェイプ名で識別して独立に検査する。
    """
    lo, hi = FOOTER_PT_RANGE
    for cn, slide in sorted(by_contract.items()):
        for shape in slide.shapes:
            if not shape.has_text_frame or shape.name not in ("Footer", "SlideNumber"):
                continue
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    if run.font.size is None or not run.text.strip():
                        continue
                    pt = run.font.size.pt
                    check(
                        lo <= pt <= hi,
                        f"{label} S{cn}: {shape.name} が {pt}pt — "
                        f"{lo:g}〜{hi:g}pt の範囲外 '{run.text[:30]}'",
                    )


def _slide_text(slide) -> str:
    """スライド内の全テキストフレームの文字列を改行区切りで連結する。"""
    return "\n".join(
        sh.text_frame.text for sh in slide.shapes if sh.has_text_frame
    )


def check_required_strings(label: str, by_contract: dict[int, object]) -> None:
    """各スライドの投影文字列が、内容契約の `Projected body` から転記した
    `REQUIRED_STRINGS` の全文字列を部分文字列として含むことを確認する。
    """
    for cn, needles in REQUIRED_STRINGS.items():
        slide = by_contract.get(cn)
        if slide is None:
            continue  # このデッキに含まれないスライド（再訪なし版の S9）
        text = _slide_text(slide)
        for needle in needles:
            check(needle in text, f"{label} S{cn}: 必須文字列が見つからない — '{needle}'")


def check_forbidden(label: str, by_contract: dict[int, object]) -> None:
    """`FORBIDDEN` の禁止表現が、許容される否定形（例: 'not ... confirmed quarry
    ponds'）の範囲外で裸に使われていないことを確認する。

    否定形パターンが与えられている場合、禁止パターンの一致区間が否定形の
    一致区間に完全に包含されていれば許容する（例が S6 に実在する
    'not individually field-confirmed quarry ponds'）。それ以外の一致は禁止表現の
    使用として検査失敗にする。
    """
    for cn, slide in sorted(by_contract.items()):
        text = _slide_text(slide)
        for positive, allowed, desc in FORBIDDEN:
            allowed_spans = [m.span() for m in re.finditer(allowed, text)] if allowed else []
            bad = [
                m for m in re.finditer(positive, text)
                if not any(a0 <= m.start() and m.end() <= a1 for a0, a1 in allowed_spans)
            ]
            check(
                not bad,
                f"{label} S{cn}: 禁止表現 '{bad[0].group(0) if bad else ''}' — {desc}",
            )


def check_english_only(label: str, by_contract: dict[int, object]) -> None:
    """投影面のテキストに CJK（漢字・かな・全角約物）が混入していないことを確認する。

    投影面は英語のみという内容契約の制約は、これまで機械検査されていなかった
    （日本語キャプションが焼き込まれた図版の混入は過去に実際に起きている）。
    ノートペインは英日併記が正であるため検査対象にしない。
    """
    for cn, slide in sorted(by_contract.items()):
        text = _slide_text(slide)
        found = CJK_PATTERN.findall(text)
        check(
            not found,
            f"{label} S{cn}: 投影テキストに CJK 文字がある — {''.join(sorted(set(found)))!r}",
        )


def check_evidence_hierarchy(label: str, by_contract: dict[int, object]) -> None:
    """evidence 階層を検査する。

    - S5: 式（`(Green − NIR)` 等）は投影本文のテキストフレームに置かず、図版内に
      置く（この検査は投影本文に存在しないことのみを確認する）。
    - S6: 春季タイル（`113`）は主結果（`145`）より小さい文字サイズであることを
      確認する。
    """
    s5 = by_contract.get(5)
    if s5 is not None:
        s5_text = _slide_text(s5)
        for formula in ("(Green − NIR)", "(Green + NIR)", "(Green − SWIR)"):
            check(
                formula not in s5_text,
                f"{label} S5: 式 '{formula}' が投影本文にある（図版内へ置く）",
            )

    def max_font_pt(slide, needle: str) -> float:
        sizes = [
            run.font.size.pt
            for sh in slide.shapes if sh.has_text_frame
            for para in sh.text_frame.paragraphs for run in para.runs
            if needle in run.text and run.font.size is not None
        ]
        return max(sizes) if sizes else 0.0

    s6 = by_contract.get(6)
    if s6 is not None:
        check(
            max_font_pt(s6, "145") > max_font_pt(s6, "113"),
            f"{label} S6: 春季113が主結果145以上の大きさになっている",
        )


def check_image_counts(label: str, by_contract: dict[int, object]) -> None:
    """各スライドの画像（PICTURE シェイプ）数が `EXPECTED_IMAGES` と一致することを確認する。

    未記載のスライド（S10）は期待値0（意図的にテキストのみ・余白を広く取る）。
    """
    for cn, slide in sorted(by_contract.items()):
        n_pic = sum(1 for sh in slide.shapes if sh.shape_type == MSO_SHAPE_TYPE.PICTURE)
        expected = EXPECTED_IMAGES.get(cn, 0)
        check(n_pic == expected, f"{label} S{cn}: 画像数 {n_pic} が期待値 {expected} と不一致")


def check_callout_range(label: str, by_contract: dict[int, object]) -> None:
    """S6 の `145`、S8 の `5–6`・`145` の callout が 60〜72pt の範囲にあることを確認する。

    Task 3 レビューの Minor 指摘: この範囲を固定する検査がなく、将来の編集で
    callout が意図せず縮小・拡大されても検査が通り続けてしまう問題への対処。
    """
    for cn, needles in CALLOUT_TARGETS.items():
        slide = by_contract.get(cn)
        if slide is None:
            continue
        for needle in needles:
            sizes = [
                run.font.size.pt
                for sh in slide.shapes if sh.has_text_frame
                for para in sh.text_frame.paragraphs for run in para.runs
                if run.text.strip() == needle and run.font.size is not None
            ]
            check(bool(sizes), f"{label} S{cn}: callout '{needle}' の文字サイズが見つからない")
            for pt in sizes:
                check(
                    CALLOUT_MIN_PT <= pt <= CALLOUT_MAX_PT,
                    f"{label} S{cn}: callout '{needle}' が {pt}pt — "
                    f"{CALLOUT_MIN_PT:.0f}〜{CALLOUT_MAX_PT:.0f}pt の範囲外",
                )


def check_pinned_photo_sources(label: str, by_contract: dict[int, object]) -> None:
    """`PINNED_PHOTO_SOURCES` に列挙したスライドの PICTURE シェイプが、意図した
    ファイルとバイト単位で一致することを確認する。

    シェイプ名やファイル名の一致だけでは、グレースケールの記事図版
    （`fig03_keirin_cliff.jpg` 等、日本語記事用に変換されたもの）や日本語キャプション
    付き図版（`fig09_multiscale.png`）への意図しないフォールバック・混入を検出
    できない（Fix round 1・Fix round 2 のレビュー指摘の再発防止）。
    """
    for cn, expected_names in PINNED_PHOTO_SOURCES.items():
        slide = by_contract.get(cn)
        if slide is None:
            continue
        pics = [sh for sh in slide.shapes if sh.shape_type == MSO_SHAPE_TYPE.PICTURE]
        check(
            len(pics) == len(expected_names),
            f"{label} S{cn}: PICTUREシェイプが{len(expected_names)}個でない: {len(pics)}個",
        )
        if len(pics) != len(expected_names):
            continue
        for pic, name in zip(pics, expected_names):
            expected_path = IMAGES / name
            check(expected_path.is_file(), f"{name} が存在しない")
            if not expected_path.is_file():
                continue
            check(
                pic.image.blob == expected_path.read_bytes(),
                f"{label} S{cn}: 画像が {name} と一致しない（バイト比較。意図しない差し替え・"
                "フォールバックの可能性）",
            )


def _images_by_blob() -> dict[bytes, str]:
    """`images/` 配下のファイル内容 → ファイル名の対応表。

    PPTX 内の画像パートはファイル名を保持しない（`image1.png` 等に付け替えられる）ため、
    配置画像がどの図版なのかはバイト比較でしか特定できない
    （`check_pinned_photo_sources` と同じ手法）。
    """
    return {p.read_bytes(): p.name for p in sorted(IMAGES.iterdir()) if p.is_file()}


def check_placed_font_sizes(label: str, by_contract: dict[int, object]) -> None:
    """図版内に焼き込まれた文字の**スライド上の実効サイズ**が15pt下限を満たすことを確認する。

    実効サイズ = native_pt × 配置倍率、配置倍率 = 配置幅 ÷ 画像実寸幅
    （実寸幅 = ピクセル幅 ÷ 画像に記録された dpi）。配置幅は生成済みPPTXの
    `shape.width` から読むので、レイアウトを変えた瞬間に検査値が追随する。

    以前は図版生成スクリプトが「配置幅220mm」という**デッキが実際には使っていない値**で
    自己申告していたため、実効 9.2〜14.7pt の図版が下限を満たしていると誤って報告して
    いた（Final review の Critical 指摘）。native サイズの正本は図版生成スクリプトの
    `NATIVE_FONT_SIZES` で、ここには複製しない。
    """
    blob_to_name = _images_by_blob()
    exercised: set[str] = set()
    for cn, slide in sorted(by_contract.items()):
        for sh in slide.shapes:
            if sh.shape_type != MSO_SHAPE_TYPE.PICTURE:
                continue
            name = blob_to_name.get(sh.image.blob)
            check(
                name is not None,
                f"{label} S{cn}: 配置画像が images/ のどのファイルとも一致しない",
            )
            if name is None or name not in NATIVE_FONT_SIZES:
                continue  # 写真は焼き込み文字を持たないため対象外
            exercised.add(name)
            px_w = sh.image.size[0]
            dpi_x = sh.image.dpi[0]
            check(dpi_x > 0, f"{label} S{cn} {name}: dpi が取得できない")
            if not dpi_x:
                continue
            scale = sh.width.inches / (px_w / dpi_x)
            for element, native_pt in sorted(NATIVE_FONT_SIZES[name].items()):
                slide_pt = native_pt * scale
                check(
                    slide_pt >= SLIDE_PT_FLOOR,
                    f"{label} S{cn} {name}: '{element}' の実効サイズ {slide_pt:.2f}pt が "
                    f"下限 {SLIDE_PT_FLOOR:g}pt 未満（native {native_pt:g}pt × "
                    f"配置倍率 {scale:.3f}）",
                )

    check(
        exercised == set(NATIVE_FONT_SIZES),
        f"{label}: native フォントサイズを宣言した図版のうち "
        f"{sorted(set(NATIVE_FONT_SIZES) - exercised)} がデッキに配置されていない",
    )


def check_placed_resolution(label: str, by_contract: dict[int, object]) -> None:
    """配置された図版の実効解像度（配置幅あたりのピクセル数）が下限を満たすことを確認する。

    実効dpi = ピクセル幅 ÷ 配置幅(インチ)。図版を小さく置けば実効dpiは上がり、
    大きく置けば下がる。`MIN_PLACED_DPI` は投影・印刷で粗さが見えない目安。
    """
    blob_to_name = _images_by_blob()
    for cn, slide in sorted(by_contract.items()):
        for sh in slide.shapes:
            if sh.shape_type != MSO_SHAPE_TYPE.PICTURE:
                continue
            name = blob_to_name.get(sh.image.blob)
            if name is None or name not in NATIVE_FONT_SIZES:
                continue  # 写真スロットは元写真の画素数で決まるため対象外
            placed_in = sh.width.inches
            placed_dpi = sh.image.size[0] / placed_in if placed_in else 0.0
            check(
                placed_dpi >= MIN_PLACED_DPI,
                f"{label} S{cn} {name}: 実効解像度 {placed_dpi:.0f}dpi が "
                f"下限 {MIN_PLACED_DPI:.0f}dpi 未満（{sh.image.size[0]}px / {placed_in:.2f}in）",
            )


def check_no_decoration_effects(label: str, path: Path) -> None:
    """作図したスライドXMLに影・グラデーション（`outerShdw` / `gradFill`）が無いことを確認する。

    検査対象は `ppt/slides/*.xml` のみ。python-pptx の既定テーマ
    （`ppt/theme/theme1.xml`）は fmtScheme にこれらを含むため、パッケージ全体を
    対象にすると「作図していない装飾」で失敗してしまう。
    """
    if not path.is_file():
        check(False, f"{label}: PPTX が存在しない（装飾効果を検査不可）")
        return
    with zipfile.ZipFile(path) as z:
        names = [n for n in z.namelist() if re.fullmatch(r"ppt/slides/slide\d+\.xml", n)]
        check(bool(names), f"{label}: ppt/slides/*.xml が見つからない")
        for member in sorted(names):
            xml = z.read(member).decode("utf-8")
            for effect in FORBIDDEN_XML_EFFECTS:
                check(
                    effect not in xml,
                    f"{label} {member}: 装飾効果 '{effect}' が使われている",
                )


def check_font_floor(label: str, by_contract: dict[int, object]) -> None:
    """本文の文字サイズが15pt未満に自動縮小されていないことを確認する。

    数字のみのラン（スライド番号や callout の裸の数値）と 12pt 以下のラン
    （footer・スライド番号として意図されたもの。範囲は `check_footer_sizes` が
    シェイプ名で独立に確認する）は例外として許容する。
    """
    for cn, slide in sorted(by_contract.items()):
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    if run.font.size is None or not run.text.strip():
                        continue
                    pt = run.font.size.pt
                    is_number_or_footer = run.text.strip().isdigit() or pt <= 12
                    check(
                        pt >= BODY_MIN_PT or is_number_or_footer,
                        f"{label} S{cn}: 本文 {pt}pt が {BODY_MIN_PT}pt 未満 — '{run.text[:30]}'",
                    )


def check_no_revisit_variant() -> None:
    """再訪なし版（11枚）から S9 が確実に抜けていることを確認する。

    スライド数・タイトル・必須文字列・禁止表現・15pt下限・画像数・callout・
    pin した画像・実効フォントサイズは、12枚版と同じ検査群を両デッキに対して
    走らせている（`main()` の `DECKS`）。ここでは「再訪スライド固有の文字列が
    残っていない」という、再訪なし版にしか意味のない不変条件だけを見る。
    """
    if not NO_REVISIT_PPTX.is_file():
        check(False, "再訪なし版のPPTXが存在しない")
        return
    prs2 = Presentation(NO_REVISIT_PPTX)
    t2 = "\n".join(sh.text_frame.text for s in prs2.slides for sh in s.shapes if sh.has_text_frame)
    check("2026-08-31" not in t2, "再訪なし版に再訪スライドが残っている")


def parse_notes(md: str) -> dict[int, dict[str, str]]:
    """スピーカーノート Markdown を `### Slide N — <title>` 単位に分解する。

    返り値は内容契約のスライド番号をキーとし、`title`（見出しのタイトル文字列）、
    `en`（英語の発話本文。`**EN (spoken)**` から `**JA (not spoken)**` の直前まで）、
    `ja`（`**JA (not spoken)**` の行から節末まで。節区切りの水平線 `---` は除く）を
    持つ辞書。

    生成スクリプト側にも同等の分解処理があるが、検査側は独立実装で持つ
    （生成スクリプトの読み取り結果をそのまま信じると、PPTX ノートペインとの
    同期検査が「自分が書いたものを読み直した」だけの検査になってしまう）。
    """
    parts = re.split(r"(?m)^### Slide (\d+) — (.*)$", md)[1:]
    sections: dict[int, dict[str, str]] = {}
    for i in range(0, len(parts), 3):
        number = int(parts[i])
        title = parts[i + 1].strip()
        body = parts[i + 2]
        en, ja = "", ""
        if EN_MARKER in body and JA_MARKER in body:
            after_en = body.split(EN_MARKER, 1)[1]
            en = after_en.split(JA_MARKER, 1)[0].strip()
            ja = re.sub(
                r"\n+-{3,}\s*$",
                "",
                (JA_MARKER + after_en.split(JA_MARKER, 1)[1]).strip(),
            ).strip()
        sections[number] = {"title": title, "en": en, "ja": ja}
    return sections


def notes_pane_text(section: dict[str, str]) -> str:
    """ノート Markdown の1節から、PPTX ノートペインに入るべき文字列を組み立てる。

    英語の発話本文を先頭に置き、空行を挟んで日本語（非発話）の塊を続ける
    （生成スクリプトの `build_notes_text()` と同一の正規形）。
    """
    return f"{section['en']}\n\n{section['ja']}"


def count_en_words(en: str) -> int:
    """英語発話本文の語数を数える。

    数字のみのトークン（`145` 等）は語として数えない。初出を綴り字
    （`one hundred and forty-five`）にする方針と整合させるため、綴られた語のみを
    通し読み時間の見積り対象とする。
    """
    return len(re.findall(r"[A-Za-z][A-Za-z'’-]*", en))


def check_speaker_notes() -> dict[int, dict[str, str]]:
    """スピーカーノート Markdown の構造・語数・S6 必須発話内容を検査する。"""
    if not NOTES_MD.is_file():
        check(False, "スピーカーノート Markdown が存在しない")
        return {}

    sections = parse_notes(NOTES_MD.read_text(encoding="utf-8"))
    check(len(sections) == 12, f"ノートのスライド数が12でない: {len(sections)}")

    for n, expected_title in enumerate(TITLES, start=1):
        label = f"S{n}"
        if n not in sections:
            check(False, f"{label}: ノートの節が存在しない")
            continue
        section = sections[n]
        check(
            section["title"] == expected_title,
            f"{label}: ノートの見出しタイトルが内容契約と完全一致しない",
        )
        check(bool(section["en"]), f"{label}: EN 発話本文が無い（{EN_MARKER}）")
        check(bool(section["ja"]), f"{label}: JA 非発話部分が無い（{JA_MARKER}）")
        check(
            EN_MARKER not in section["en"] and JA_MARKER not in section["en"],
            f"{label}: EN 部分に構造マーカーが混入している",
        )
        words = count_en_words(section["en"])
        budget = round(DURATIONS_S[n - 1] / 60 * WPM)
        check(
            abs(words - budget) <= budget * WORD_TOLERANCE,
            f"{label}: EN 語数 {words} が想定 {budget} 語から "
            f"{WORD_TOLERANCE:.0%} 超乖離",
        )

    en6 = sections.get(6, {}).get("en", "")
    for description, pattern in S6_REQUIRED_FACTS:
        check(
            re.search(pattern, en6) is not None,
            f"S6: required spoken content の事実が EN 本文に無い — {description}",
        )

    check_revisit_update_marker(sections)
    return sections


def revisit_photos_are_placeholders() -> bool:
    """S9 の写真スロットが実写に差し替わっていないかを判定する。

    生成スクリプトは `images/revisit_{n}.jpg`（`.jpeg` / `.png` も可）があればそれを、
    無ければプレースホルダを使う。2枚のうち片方でも実写が無ければプレースホルダ運用中と
    みなす（解決順を生成側と一致させないと、`.png` で納品された場合に判定がずれる）。
    """
    return not all(
        any((IMAGES / f"{stem}.{ext}").is_file() for ext in REVISIT_PHOTO_EXTS)
        for stem in REVISIT_PHOTO_STEMS
    )


def check_revisit_update_marker(sections: dict[int, dict[str, str]]) -> None:
    """S9 の再訪ノートの見直し要求が、写真がプレースホルダの間だけ立つことを確認する。

    S9 の英語本文は 8/31 の再訪より前に書いたものであり、まだ起きていない出来事を
    語らない構成にしてある。写真が実写に差し替わったら本文も見直す必要があるため、
    プレースホルダ運用中は `REVISIT_UPDATE_MARKER` を **JA（非発話）側**に置くことを
    必須とする。実写2枚が揃った時点でこの要求は自動的に消える。

    目印は英語の発話本文に入ってはならない（英語は読み上げるため）。この禁止は
    写真の有無にかかわらず常に検査する。
    """
    en9 = sections.get(9, {}).get("en", "")
    check(
        REVISIT_UPDATE_MARKER not in en9,
        f"S9: '{REVISIT_UPDATE_MARKER}' が EN 発話本文にある（読み上げてしまう）",
    )
    if not revisit_photos_are_placeholders():
        return
    ja9 = sections.get(9, {}).get("ja", "")
    check(
        REVISIT_UPDATE_MARKER in ja9,
        f"S9: 再訪写真がプレースホルダのままなのに JA 側に "
        f"'{REVISIT_UPDATE_MARKER}' が無い（再訪後の本文見直しの申し送りが消えている）",
    )


def check_notes_sync(sections: dict[int, dict[str, str]]) -> None:
    """両版の PPTX ノートペインが、ノート Markdown の正規形と一致することを確認する。

    12枚版はスライド位置＝内容契約の番号だが、再訪なし版（11枚）は S9 が抜けるため
    `NO_REVISIT_CONTRACT_NUMBERS` で対応付ける（位置で素朴に添字を取ると、
    位置9以降に S10・S11・S12 ではなく1つ前のノートが載っていても検査が通ってしまう）。
    """
    if not sections:
        return
    decks = [
        ("12枚版", PPTX, list(range(1, 13))),
        ("再訪なし版", NO_REVISIT_PPTX, NO_REVISIT_CONTRACT_NUMBERS),
    ]
    for deck_label, path, contract_numbers in decks:
        if not path.is_file():
            check(False, f"{deck_label}: PPTX が存在しない（ノート同期を検査不可）")
            continue
        slides = list(Presentation(path).slides)
        check(
            len(slides) == len(contract_numbers),
            f"{deck_label}: スライド数 {len(slides)} が期待値 {len(contract_numbers)} と不一致",
        )
        for position, contract_n in enumerate(contract_numbers, start=1):
            label = f"{deck_label} 位置{position}（契約 S{contract_n}）"
            if position > len(slides):
                check(False, f"{label}: スライドが存在しない")
                continue
            slide = slides[position - 1]
            check(slide.has_notes_slide, f"{label}: ノートペインが無い")
            if not slide.has_notes_slide:
                continue
            actual = slide.notes_slide.notes_text_frame.text
            expected = notes_pane_text(sections[contract_n])
            check(
                actual == expected,
                f"{label}: PPTX ノートペインが Markdown と一致しない",
            )


HEX64 = re.compile(r"\b[0-9a-f]{64}\b")

# Task 6 照合記録が出典を示す義務を負う数値。口頭・投影のどちらかで使われる値の
# 全体集合（ブリーフ Step 1 の最小集合 + タスク本体が追加する 20 m・負の閾値等）。
# 検出に幅を持たせるため、裸の数字（"0.3" 等）は誤検出を避けられる形にアンカーする。
VERIFICATION_NUMBER_VALUES: tuple[str, ...] = (
    "145", "113", "127", "9 px", "100 m²", "10 m", "20 m", "1.28 ha",
    "7,826 m²", "2025-03-23", "2025-08-02", "0.0%", "0.7%",
    "−0.2", "−0.1", "NDVI > 0.3",
)

# 照合記録に名前だけ残すべき、削除済み・不採用の図版（ブリーフ Step 1 が明示的に
# 名前を要求する）。images/ には存在しないため、SHA256照合の対象からは除外する。
VERIFICATION_RETIRED_IMAGE_NAMES: tuple[str, ...] = (
    "fig03_keirin_cliff.jpg", "poster_f4_index_panels.png",
    # 2026-08-24 の調整で S4・S7パネル(b) を色付き原本へ差し替えた結果、どのスライドからも
    # 参照されなくなり `images/` から削除した2点。削除の経緯を記録に残すため名前で検査する
    # （`docs/articles/` 側の原本は無変更で存在する）。
    "fig05_drone_takeoff.jpg", "fig06_aerial_quarries.jpg",
    # 2026-08-24 の是正で背景地図を航空写真（地理院タイル `seamlessphoto`）から
    # ラベルなしの淡色地図（CARTO Positron）へ差し替えた結果、参照されなくなり
    # `images/` から削除したラスタ。差し替えの経緯（S7 の三スケールの物語との衝突と、
    # ground truth と読まれる懸念）を記録に残すため名前で検査する。
    "basemap_kitagi_gsi_seamlessphoto.png",
)

# 8/31 撮影分の空欄テーブルに必須の見出し4列。
VERIFICATION_PHOTO_TABLE_HEADERS: tuple[str, ...] = (
    "座標", "撮影時刻", "撮影方向", "対象ポリゴンID",
)


def check_verification_record() -> None:
    """Task 6 照合記録（数値・出典・画像のSHA256・8/31写真の空欄表）を検査する。

    `poster_f4_index_panels.png` は現在も `images/` に存在するため、SHA256照合
    （`images/` 全ファイルのループ）の対象にも自然に入る。ブリーフはこれを
    `fig03_keirin_cliff.jpg`（削除済み）と並べて名前検査の対象に挙げているため、
    ここでは両方を独立に検査する。
    """
    if not VERIFICATION.is_file():
        check(False, "照合記録が存在しない")
        return
    v = VERIFICATION.read_text(encoding="utf-8")

    for value in VERIFICATION_NUMBER_VALUES:
        check(value in v, f"照合記録に値の出典が無い: {value}")

    check("SHA256" in v, "照合記録に SHA256 の記載が無い")

    for name in VERIFICATION_RETIRED_IMAGE_NAMES:
        check(name in v, f"照合記録に画像の記載が無い: {name}")

    # images/ 配下の全ファイルが記録に載っており、記録中のSHA256がディスク上の
    # ファイルを再計算したハッシュと一致することを確認する（スクリプトが自ら
    # 再計算する。過去のレポートからのコピーでは古いハッシュの混入を検出できない）。
    image_files = sorted(p for p in IMAGES.iterdir() if p.is_file())
    check(bool(image_files), "images/ にファイルが無い（検査不可）")
    for path in image_files:
        name = path.name
        check(name in v, f"照合記録に画像の記載が無い: {name}")
        actual_hash = hashlib.sha256(path.read_bytes()).hexdigest()
        row_hashes = [
            hexs
            for line in v.splitlines()
            if name in line
            for hexs in HEX64.findall(line)
        ]
        check(bool(row_hashes), f"{name}: 照合記録にSHA256の記載行が無い")
        check(
            actual_hash in row_hashes,
            f"{name}: 記録のSHA256がディスク上のファイルと一致しない"
            f"（記録: {row_hashes}, 実際: {actual_hash}）",
        )

    for header in VERIFICATION_PHOTO_TABLE_HEADERS:
        check(header in v, f"照合記録に8/31撮影テーブルの見出しが無い: {header}")


# 検査を走らせる2つのデッキ。値は「デッキ内の各位置に対応する内容契約のスライド番号」。
# 12枚版は位置＝契約番号、再訪なし版は S9 が抜けて位置9以降がずれる。
DECKS: list[tuple[str, Path, list[int]]] = [
    ("12枚版", PPTX, list(range(1, 13))),
    ("再訪なし版", NO_REVISIT_PPTX, NO_REVISIT_CONTRACT_NUMBERS),
]


def main() -> None:
    check_figures()

    # 投影面に関する検査群は両デッキに同じものを走らせる。以前はタイトル・必須文字列・
    # 禁止表現・15pt下限・画像数・callout・pin した画像がすべて12枚版だけを見ており、
    # 再訪なし版はスライド数と1文字列しか検査されていなかった（Final review の指摘）。
    for label, path, contract_numbers in DECKS:
        by_contract = deck_slides(label, path, contract_numbers)
        if by_contract is None:
            continue
        check_titles(label, by_contract)
        check_required_strings(label, by_contract)
        check_forbidden(label, by_contract)
        check_english_only(label, by_contract)
        check_evidence_hierarchy(label, by_contract)
        check_font_floor(label, by_contract)
        check_footer_sizes(label, by_contract)
        check_image_counts(label, by_contract)
        check_callout_range(label, by_contract)
        check_pinned_photo_sources(label, by_contract)
        check_placed_font_sizes(label, by_contract)
        check_placed_resolution(label, by_contract)
        check_no_decoration_effects(label, path)

    check_no_revisit_variant()

    sections = check_speaker_notes()
    check_notes_sync(sections)

    check_verification_record()

    if errors:
        print(f"FAIL ({len(errors)} / {checks} checks failed)")
        for e in errors:
            print(f"  - {e}")
        sys.exit(1)
    print(f"OK: {checks} checks passed")


if __name__ == "__main__":
    main()
